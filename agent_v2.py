#!/usr/bin/env python3
"""
Jane's AI Agent v2 - Modular Local AI Assistant
Works with Ollama + local models (qwen2.5-coder:14b)
No cloud APIs needed.

Usage:
    python agent_v2.py

Modules:
    - File Manager: organize, clean, create files/folders
    - Web Parser: scrape websites, take screenshots
    - Document Generator: create Word, PDF, TXT, Excel files
    - System Commands: open apps, browser, manage PC
    - Email (Gmail): read, analyze, send emails
    - Search & Analyze: web search, summarize info
"""

import json
import subprocess
import requests
import os
import sys
import shutil
import glob
import datetime
import re
from pathlib import Path

# ============================================================
# CONFIGURATION - Edit these settings
# ============================================================
OLLAMA_URL = "http://localhost:11434"
MODEL = "qwen2.5-coder:14b"

# Windows paths
DESKTOP = os.path.join(os.environ.get('USERPROFILE', os.path.expanduser('~')), 'Desktop')
DOCUMENTS = os.path.join(os.environ.get('USERPROFILE', os.path.expanduser('~')), 'Documents')
DOWNLOADS = os.path.join(os.environ.get('USERPROFILE', os.path.expanduser('~')), 'Downloads')

# Agent working directory (for logs, temp files)
AGENT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGS_DIR = os.path.join(AGENT_DIR, "logs")
os.makedirs(LOGS_DIR, exist_ok=True)

# ============================================================
# AVAILABLE TOOLS DESCRIPTION (sent to LLM)
# ============================================================
TOOLS_DESCRIPTION = """
You have access to these tools. To use a tool, respond with JSON:
{"tool": "tool_name", "args": {"arg1": "value1", "arg2": "value2"}}

AVAILABLE TOOLS:

1. run_python - Execute Python code
   {"tool": "run_python", "args": {"code": "print('hello')"}}

2. run_powershell - Execute PowerShell command
   {"tool": "run_powershell", "args": {"code": "Get-ChildItem"}}

3. create_file - Create a text file with content
   {"tool": "create_file", "args": {"path": "C:/Users/Dator/Desktop/test.txt", "content": "Hello world"}}

4. read_file - Read file contents
   {"tool": "read_file", "args": {"path": "C:/Users/Dator/Desktop/test.txt"}}

5. list_files - List files in a directory (with optional filter)
   {"tool": "list_files", "args": {"path": "C:/Users/Dator/Desktop", "pattern": "*.txt"}}

6. organize_folder - Auto-organize files in folder by type
   {"tool": "organize_folder", "args": {"path": "C:/Users/Dator/Downloads"}}

7. find_duplicates - Find duplicate files in a folder
   {"tool": "find_duplicates", "args": {"path": "C:/Users/Dator/Documents"}}

8. disk_usage - Show disk space and large files
   {"tool": "disk_usage", "args": {"path": "C:/"}}

9. open_app - Open an application or URL
   {"tool": "open_app", "args": {"target": "https://mail.google.com"}}
   {"tool": "open_app", "args": {"target": "notepad"}}

10. create_document - Create Word/PDF/Excel document
    {"tool": "create_document", "args": {"type": "word", "path": "C:/Users/Dator/Desktop/report.docx", "title": "Report Title", "content": "Document text here"}}
    types: word, pdf, excel, presentation

11. web_search - Search the web (via DuckDuckGo)
    {"tool": "web_search", "args": {"query": "best investments 2025"}}

12. parse_webpage - Parse a webpage and extract text
    {"tool": "parse_webpage", "args": {"url": "https://example.com", "selector": "article"}}

13. take_screenshot - Screenshot a webpage
    {"tool": "take_screenshot", "args": {"url": "https://example.com", "path": "C:/Users/Dator/Desktop/shot.png"}}

14. system_info - Get system information (CPU, RAM, disk)
    {"tool": "system_info", "args": {}}

15. clean_temp - Clean temporary files and caches
    {"tool": "clean_temp", "args": {"dry_run": true}}

16. send_email - Send email via Gmail (needs app password setup)
    {"tool": "send_email", "args": {"to": "someone@email.com", "subject": "Hello", "body": "Message text"}}

17. chat - Just respond with text, no action needed
    {"tool": "chat", "args": {"message": "Your response here"}}

IMPORTANT RULES:
- ALWAYS respond with ONE JSON object. Nothing else.
- Use encoding='utf-8' when writing files with non-English text.
- Use forward slashes (/) in paths, or raw strings.
- For Windows user folder: C:/Users/Dator/
- Desktop: C:/Users/Dator/Desktop/
- Documents: C:/Users/Dator/Documents/
- Downloads: C:/Users/Dator/Downloads/
"""

SYSTEM_PROMPT = f"""You are Jane's personal AI assistant running on her Windows PC.
You help with file management, web parsing, document creation, system management, and more.
You speak Russian and English.

{TOOLS_DESCRIPTION}

Respond ONLY with a JSON object. No explanations before or after."""


# ============================================================
# TOOL IMPLEMENTATIONS
# ============================================================

def tool_run_python(code: str) -> str:
    """Execute Python code"""
    try:
        result = subprocess.run(
            [sys.executable, "-c", code],
            capture_output=True, text=True, timeout=60,
            encoding="utf-8", errors="replace"
        )
        output = ""
        if result.stdout:
            output += result.stdout
        if result.stderr:
            output += "\n[STDERR] " + result.stderr
        return output.strip() if output.strip() else "Done. No output."
    except subprocess.TimeoutExpired:
        return "Error: command timed out (60s limit)"
    except Exception as e:
        return f"Error: {e}"


def tool_run_powershell(code: str) -> str:
    """Execute PowerShell command"""
    try:
        result = subprocess.run(
            ["powershell", "-Command", code],
            capture_output=True, text=True, timeout=60,
            encoding="utf-8", errors="replace"
        )
        output = ""
        if result.stdout:
            output += result.stdout
        if result.stderr:
            output += "\n[STDERR] " + result.stderr
        return output.strip() if output.strip() else "Done. No output."
    except subprocess.TimeoutExpired:
        return "Error: command timed out (60s limit)"
    except Exception as e:
        return f"Error: {e}"


def tool_create_file(path: str, content: str) -> str:
    """Create a file with content"""
    try:
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, 'w', encoding='utf-8') as f:
            f.write(content)
        return f"File created: {path}"
    except Exception as e:
        return f"Error creating file: {e}"


def tool_read_file(path: str) -> str:
    """Read file contents"""
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read(50000)  # Limit to 50KB
        return content if content else "(empty file)"
    except Exception as e:
        return f"Error reading file: {e}"


def tool_list_files(path: str, pattern: str = "*") -> str:
    """List files in directory"""
    try:
        files = []
        search_path = os.path.join(path, pattern)
        for f in glob.glob(search_path):
            stat = os.stat(f)
            size = stat.st_size
            modified = datetime.datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M')
            name = os.path.basename(f)
            ftype = "DIR" if os.path.isdir(f) else "FILE"
            
            # Human-readable size
            if size < 1024:
                size_str = f"{size} B"
            elif size < 1024*1024:
                size_str = f"{size/1024:.1f} KB"
            elif size < 1024*1024*1024:
                size_str = f"{size/(1024*1024):.1f} MB"
            else:
                size_str = f"{size/(1024*1024*1024):.1f} GB"
            
            files.append(f"  {ftype}  {size_str:>10}  {modified}  {name}")
        
        if not files:
            return f"No files matching '{pattern}' in {path}"
        
        header = f"Files in {path} (pattern: {pattern}):\n"
        return header + "\n".join(sorted(files)[:100])  # Limit to 100 items
    except Exception as e:
        return f"Error listing files: {e}"


def tool_organize_folder(path: str) -> str:
    """Organize files in folder by type"""
    categories = {
        'Images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.webp', '.ico', '.tiff'],
        'Documents': ['.pdf', '.doc', '.docx', '.txt', '.rtf', '.odt', '.xls', '.xlsx', '.csv', '.pptx', '.ppt'],
        'Videos': ['.mp4', '.avi', '.mkv', '.mov', '.wmv', '.flv', '.webm'],
        'Audio': ['.mp3', '.wav', '.flac', '.aac', '.ogg', '.wma'],
        'Archives': ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2'],
        'Code': ['.py', '.js', '.html', '.css', '.java', '.cpp', '.c', '.json', '.xml'],
        'Installers': ['.exe', '.msi', '.dmg', '.deb', '.rpm'],
    }
    
    try:
        moved = 0
        report = []
        
        for item in os.listdir(path):
            item_path = os.path.join(path, item)
            if os.path.isdir(item_path):
                continue
            
            ext = os.path.splitext(item)[1].lower()
            target_folder = 'Other'
            
            for category, extensions in categories.items():
                if ext in extensions:
                    target_folder = category
                    break
            
            target_dir = os.path.join(path, target_folder)
            os.makedirs(target_dir, exist_ok=True)
            
            target_path = os.path.join(target_dir, item)
            if os.path.exists(target_path):
                base, extension = os.path.splitext(item)
                target_path = os.path.join(target_dir, f"{base}_{datetime.datetime.now().strftime('%H%M%S')}{extension}")
            
            shutil.move(item_path, target_path)
            moved += 1
            report.append(f"  {item} -> {target_folder}/")
        
        if moved == 0:
            return f"No files to organize in {path}"
        
        result = f"Organized {moved} files in {path}:\n"
        result += "\n".join(report[:50])
        if len(report) > 50:
            result += f"\n  ... and {len(report) - 50} more"
        return result
    except Exception as e:
        return f"Error organizing folder: {e}"


def tool_find_duplicates(path: str) -> str:
    """Find duplicate files by size and name"""
    try:
        files_by_size = {}
        for root, dirs, files in os.walk(path):
            for name in files:
                filepath = os.path.join(root, name)
                try:
                    size = os.path.getsize(filepath)
                    key = (name, size)
                    if key not in files_by_size:
                        files_by_size[key] = []
                    files_by_size[key].append(filepath)
                except OSError:
                    continue
        
        duplicates = {k: v for k, v in files_by_size.items() if len(v) > 1}
        
        if not duplicates:
            return f"No duplicates found in {path}"
        
        result = f"Found {len(duplicates)} groups of duplicates:\n"
        for (name, size), paths in list(duplicates.items())[:20]:
            size_str = f"{size/1024:.1f} KB" if size < 1024*1024 else f"{size/(1024*1024):.1f} MB"
            result += f"\n  {name} ({size_str}):\n"
            for p in paths:
                result += f"    - {p}\n"
        
        return result
    except Exception as e:
        return f"Error finding duplicates: {e}"


def tool_disk_usage(path: str = "C:/") -> str:
    """Show disk usage info"""
    try:
        total, used, free = shutil.disk_usage(path)
        result = f"Disk {path}:\n"
        result += f"  Total: {total/(1024**3):.1f} GB\n"
        result += f"  Used:  {used/(1024**3):.1f} GB ({used/total*100:.1f}%)\n"
        result += f"  Free:  {free/(1024**3):.1f} GB ({free/total*100:.1f}%)\n"
        
        # Find large files in common locations
        result += f"\nLargest files in common folders:\n"
        search_dirs = [DOWNLOADS, DOCUMENTS, DESKTOP]
        large_files = []
        
        for search_dir in search_dirs:
            if os.path.exists(search_dir):
                for root, dirs, files in os.walk(search_dir):
                    for name in files:
                        try:
                            fp = os.path.join(root, name)
                            size = os.path.getsize(fp)
                            if size > 50 * 1024 * 1024:  # > 50MB
                                large_files.append((size, fp))
                        except OSError:
                            continue
        
        large_files.sort(reverse=True)
        for size, fp in large_files[:15]:
            result += f"  {size/(1024**2):.0f} MB  {fp}\n"
        
        if not large_files:
            result += "  No files > 50MB found in common folders\n"
        
        return result
    except Exception as e:
        return f"Error: {e}"


def tool_open_app(target: str) -> str:
    """Open application or URL"""
    try:
        if target.startswith(('http://', 'https://')):
            subprocess.Popen(["powershell", "-Command", f"Start-Process '{target}'"])
            return f"Opened in browser: {target}"
        else:
            subprocess.Popen(["powershell", "-Command", f"Start-Process {target}"])
            return f"Opened: {target}"
    except Exception as e:
        return f"Error opening {target}: {e}"


def tool_create_document(doc_type: str, path: str, title: str = "", content: str = "") -> str:
    """Create Word/PDF/Excel document"""
    try:
        if doc_type == "word":
            # Try python-docx, fallback to simple approach
            try:
                from docx import Document
                doc = Document()
                if title:
                    doc.add_heading(title, 0)
                for paragraph in content.split('\n'):
                    if paragraph.strip():
                        doc.add_paragraph(paragraph)
                doc.save(path)
                return f"Word document created: {path}"
            except ImportError:
                return "Error: python-docx not installed. Run: pip install python-docx"
        
        elif doc_type == "pdf":
            try:
                from reportlab.lib.pagesizes import A4
                from reportlab.lib.units import cm
                from reportlab.pdfgen import canvas
                from reportlab.pdfbase import pdfmetrics
                from reportlab.pdfbase.ttfonts import TTFont
                
                c = canvas.Canvas(path, pagesize=A4)
                width, height = A4
                
                # Try to register a font that supports Cyrillic
                try:
                    pdfmetrics.registerFont(TTFont('Arial', 'C:/Windows/Fonts/arial.ttf'))
                    font_name = 'Arial'
                except:
                    font_name = 'Helvetica'
                
                y = height - 2*cm
                
                if title:
                    c.setFont(font_name, 18)
                    c.drawString(2*cm, y, title)
                    y -= 1.5*cm
                
                c.setFont(font_name, 12)
                for line in content.split('\n'):
                    if y < 2*cm:
                        c.showPage()
                        y = height - 2*cm
                        c.setFont(font_name, 12)
                    c.drawString(2*cm, y, line)
                    y -= 0.5*cm
                
                c.save()
                return f"PDF created: {path}"
            except ImportError:
                return "Error: reportlab not installed. Run: pip install reportlab"
        
        elif doc_type == "excel":
            try:
                import openpyxl
                wb = openpyxl.Workbook()
                ws = wb.active
                if title:
                    ws.title = title[:31]  # Sheet name max 31 chars
                
                for i, line in enumerate(content.split('\n'), 1):
                    cells = line.split('\t') if '\t' in line else line.split(',')
                    for j, cell in enumerate(cells, 1):
                        ws.cell(row=i, column=j, value=cell.strip())
                
                wb.save(path)
                return f"Excel file created: {path}"
            except ImportError:
                return "Error: openpyxl not installed. Run: pip install openpyxl"
        
        elif doc_type == "presentation":
            try:
                from pptx import Presentation
                prs = Presentation()
                
                # Title slide
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title or "Presentation"
                
                # Content slides (split by double newline)
                sections = content.split('\n\n')
                for section in sections:
                    if section.strip():
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        lines = section.strip().split('\n')
                        slide.shapes.title.text = lines[0]
                        if len(lines) > 1:
                            body = slide.placeholders[1]
                            body.text = '\n'.join(lines[1:])
                
                prs.save(path)
                return f"Presentation created: {path}"
            except ImportError:
                return "Error: python-pptx not installed. Run: pip install python-pptx"
        
        else:
            return f"Unknown document type: {doc_type}. Use: word, pdf, excel, presentation"
    
    except Exception as e:
        return f"Error creating document: {e}"


def tool_web_search(query: str) -> str:
    """Search the web via DuckDuckGo"""
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        resp = requests.get(
            f"https://html.duckduckgo.com/html/?q={requests.utils.quote(query)}",
            headers=headers, timeout=15
        )
        
        # Simple extraction of results
        results = []
        from html.parser import HTMLParser
        
        class DDGParser(HTMLParser):
            def __init__(self):
                super().__init__()
                self.in_result = False
                self.in_title = False
                self.in_snippet = False
                self.current = {}
                self.results = []
            
            def handle_starttag(self, tag, attrs):
                attrs_dict = dict(attrs)
                if tag == 'a' and 'result__a' in attrs_dict.get('class', ''):
                    self.in_title = True
                    self.current = {'title': '', 'url': attrs_dict.get('href', ''), 'snippet': ''}
                elif tag == 'a' and 'result__snippet' in attrs_dict.get('class', ''):
                    self.in_snippet = True
            
            def handle_data(self, data):
                if self.in_title:
                    self.current['title'] += data
                elif self.in_snippet:
                    self.current['snippet'] += data
            
            def handle_endtag(self, tag):
                if tag == 'a':
                    if self.in_title:
                        self.in_title = False
                    elif self.in_snippet:
                        self.in_snippet = False
                        if self.current.get('title'):
                            self.results.append(self.current)
                            self.current = {}
        
        parser = DDGParser()
        parser.feed(resp.text)
        
        if not parser.results:
            return f"No results found for: {query}"
        
        result = f"Search results for '{query}':\n\n"
        for i, r in enumerate(parser.results[:8], 1):
            result += f"{i}. {r['title']}\n   {r['url']}\n   {r['snippet']}\n\n"
        
        return result
    except Exception as e:
        return f"Search error: {e}"


def tool_parse_webpage(url: str, selector: str = None) -> str:
    """Parse webpage and extract text"""
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        resp = requests.get(url, headers=headers, timeout=15)
        resp.encoding = resp.apparent_encoding or 'utf-8'
        
        # Simple HTML to text
        from html.parser import HTMLParser
        
        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []
                self.skip = False
                self.skip_tags = {'script', 'style', 'head', 'nav', 'footer', 'header'}
            
            def handle_starttag(self, tag, attrs):
                if tag in self.skip_tags:
                    self.skip = True
            
            def handle_endtag(self, tag):
                if tag in self.skip_tags:
                    self.skip = False
                if tag in ('p', 'div', 'br', 'h1', 'h2', 'h3', 'h4', 'li', 'tr'):
                    self.text.append('\n')
            
            def handle_data(self, data):
                if not self.skip:
                    text = data.strip()
                    if text:
                        self.text.append(text)
        
        extractor = TextExtractor()
        extractor.feed(resp.text)
        text = ' '.join(extractor.text)
        
        # Clean up whitespace
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = text[:5000]  # Limit output
        
        # Extract title
        title_match = re.search(r'<title>(.*?)</title>', resp.text, re.IGNORECASE | re.DOTALL)
        title = title_match.group(1).strip() if title_match else "No title"
        
        return f"Title: {title}\nURL: {url}\n\nContent:\n{text}"
    except Exception as e:
        return f"Error parsing {url}: {e}"


def tool_take_screenshot(url: str, path: str = None) -> str:
    """Take screenshot of webpage using Playwright"""
    if not path:
        path = os.path.join(DESKTOP, f"screenshot_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
    
    # Use Playwright if available, otherwise use PowerShell approach
    code = f"""
import sys
try:
    from playwright.sync_api import sync_playwright
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        page.goto("{url}", wait_until="domcontentloaded")
        page.screenshot(path=r"{path}", full_page=True)
        browser.close()
        print(f"Screenshot saved: {path}")
except ImportError:
    print("Playwright not installed. Run: pip install playwright && playwright install chromium")
except Exception as e:
    print(f"Error: {{e}}")
"""
    return tool_run_python(code)


def tool_system_info() -> str:
    """Get system information"""
    code = """
import platform
import psutil
import datetime

print(f"=== System Info ===")
print(f"OS: {platform.system()} {platform.release()} ({platform.version()})")
print(f"Machine: {platform.machine()}")
print(f"Processor: {platform.processor()}")
print(f"")
print(f"=== CPU ===")
print(f"Cores: {psutil.cpu_count(logical=False)} physical, {psutil.cpu_count()} logical")
print(f"Usage: {psutil.cpu_percent(interval=1)}%")
print(f"")
print(f"=== Memory ===")
mem = psutil.virtual_memory()
print(f"Total: {mem.total/(1024**3):.1f} GB")
print(f"Used: {mem.used/(1024**3):.1f} GB ({mem.percent}%)")
print(f"Free: {mem.available/(1024**3):.1f} GB")
print(f"")
print(f"=== Disk ===")
for part in psutil.disk_partitions():
    try:
        usage = psutil.disk_usage(part.mountpoint)
        print(f"{part.device}: {usage.used/(1024**3):.1f}/{usage.total/(1024**3):.1f} GB ({usage.percent}%)")
    except:
        pass
print(f"")
print(f"=== Top Processes (by memory) ===")
procs = []
for proc in psutil.process_iter(['name', 'memory_percent']):
    try:
        procs.append((proc.info['memory_percent'], proc.info['name']))
    except:
        pass
procs.sort(reverse=True)
for mem_pct, name in procs[:10]:
    print(f"  {mem_pct:.1f}%  {name}")
"""
    return tool_run_python(code)


def tool_clean_temp(dry_run: bool = True) -> str:
    """Clean temporary files"""
    temp_dirs = [
        os.environ.get('TEMP', ''),
        os.path.join(os.environ.get('LOCALAPPDATA', ''), 'Temp'),
    ]
    
    result = "=== Temp Files Cleanup ===\n"
    total_size = 0
    total_files = 0
    
    for temp_dir in temp_dirs:
        if not temp_dir or not os.path.exists(temp_dir):
            continue
        
        result += f"\n{temp_dir}:\n"
        dir_size = 0
        dir_files = 0
        
        for root, dirs, files in os.walk(temp_dir):
            for name in files:
                try:
                    fp = os.path.join(root, name)
                    size = os.path.getsize(fp)
                    dir_size += size
                    dir_files += 1
                    
                    if not dry_run:
                        try:
                            os.remove(fp)
                        except (PermissionError, OSError):
                            pass
                except OSError:
                    continue
        
        result += f"  Files: {dir_files}, Size: {dir_size/(1024**2):.1f} MB\n"
        total_size += dir_size
        total_files += dir_files
    
    mode = "DRY RUN (no files deleted)" if dry_run else "CLEANED"
    result += f"\n--- {mode} ---\n"
    result += f"Total: {total_files} files, {total_size/(1024**2):.1f} MB\n"
    
    if dry_run:
        result += "\nTo actually delete, run with dry_run=false"
    
    return result


def tool_send_email(to: str, subject: str, body: str) -> str:
    """Send email via Gmail (requires app password in .env)"""
    env_file = os.path.join(AGENT_DIR, ".env")
    
    if not os.path.exists(env_file):
        return """Email not configured yet. To set up Gmail:
1. Go to https://myaccount.google.com/apppasswords
2. Create an App Password for 'Mail'
3. Create file .env next to agent_v2.py with:
   GMAIL_ADDRESS=your.email@gmail.com
   GMAIL_APP_PASSWORD=your_app_password
4. Try again"""
    
    # Read .env
    gmail_address = ""
    gmail_password = ""
    with open(env_file, 'r') as f:
        for line in f:
            if line.startswith('GMAIL_ADDRESS='):
                gmail_address = line.split('=', 1)[1].strip()
            elif line.startswith('GMAIL_APP_PASSWORD='):
                gmail_password = line.split('=', 1)[1].strip()
    
    if not gmail_address or not gmail_password:
        return "Error: GMAIL_ADDRESS or GMAIL_APP_PASSWORD not found in .env file"
    
    try:
        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart
        
        msg = MIMEMultipart()
        msg['From'] = gmail_address
        msg['To'] = to
        msg['Subject'] = subject
        msg.attach(MIMEText(body, 'plain', 'utf-8'))
        
        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
            server.login(gmail_address, gmail_password)
            server.send_message(msg)
        
        return f"Email sent to {to} with subject '{subject}'"
    except Exception as e:
        return f"Error sending email: {e}"


# ============================================================
# TOOL ROUTER
# ============================================================

TOOLS = {
    "run_python": lambda args: tool_run_python(args.get("code", "")),
    "run_powershell": lambda args: tool_run_powershell(args.get("code", "")),
    "create_file": lambda args: tool_create_file(args.get("path", ""), args.get("content", "")),
    "read_file": lambda args: tool_read_file(args.get("path", "")),
    "list_files": lambda args: tool_list_files(args.get("path", ""), args.get("pattern", "*")),
    "organize_folder": lambda args: tool_organize_folder(args.get("path", "")),
    "find_duplicates": lambda args: tool_find_duplicates(args.get("path", "")),
    "disk_usage": lambda args: tool_disk_usage(args.get("path", "C:/")),
    "open_app": lambda args: tool_open_app(args.get("target", "")),
    "create_document": lambda args: tool_create_document(
        args.get("type", "word"), args.get("path", ""),
        args.get("title", ""), args.get("content", "")
    ),
    "web_search": lambda args: tool_web_search(args.get("query", "")),
    "parse_webpage": lambda args: tool_parse_webpage(args.get("url", ""), args.get("selector")),
    "take_screenshot": lambda args: tool_take_screenshot(args.get("url", ""), args.get("path")),
    "system_info": lambda args: tool_system_info(),
    "clean_temp": lambda args: tool_clean_temp(args.get("dry_run", True)),
    "send_email": lambda args: tool_send_email(args.get("to", ""), args.get("subject", ""), args.get("body", "")),
    "chat": lambda args: args.get("message", ""),
}


# ============================================================
# LLM COMMUNICATION
# ============================================================

def ask_ollama(user_message: str, history: list) -> str:
    """Send message to Ollama and get response"""
    history.append({"role": "user", "content": user_message})
    
    messages = [{"role": "system", "content": SYSTEM_PROMPT}] + history[-20:]  # Keep last 20 messages
    
    try:
        resp = requests.post(f"{OLLAMA_URL}/api/chat", json={
            "model": MODEL,
            "messages": messages,
            "stream": False,
            "options": {"temperature": 0.3, "num_predict": 4096}
        }, timeout=120)
        
        answer = resp.json()["message"]["content"]
        history.append({"role": "assistant", "content": answer})
        return answer
    except requests.exceptions.ConnectionError:
        return '{"tool": "chat", "args": {"message": "Error: Cannot connect to Ollama. Is it running? Start Ollama first."}}'
    except Exception as e:
        return f'{{"tool": "chat", "args": {{"message": "Error communicating with Ollama: {e}"}}}}'


def extract_json(text: str) -> dict:
    """Extract JSON from model response"""
    text = text.strip()
    
    # Remove thinking tags
    if "<think>" in text:
        idx = text.rfind("</think>")
        if idx != -1:
            text = text[idx + 8:].strip()
    
    # Remove markdown code blocks
    if "```json" in text:
        text = text.split("```json")[1].split("```")[0].strip()
    elif "```" in text:
        text = text.split("```")[1].split("```")[0].strip()
    
    # Find JSON
    start = text.find("{")
    end = text.rfind("}") + 1
    if start != -1 and end > start:
        try:
            return json.loads(text[start:end])
        except json.JSONDecodeError:
            # Try to fix common JSON issues
            json_str = text[start:end]
            json_str = json_str.replace("'", '"')
            try:
                return json.loads(json_str)
            except:
                pass
    
    return None


def guess_tool(text: str) -> dict:
    """Fallback: try to detect a tool call from plain text"""
    text = text.strip()
    lower = text.lower()
    
    # Remove thinking
    if "<think>" in text:
        idx = text.rfind("</think>")
        if idx != -1:
            text = text[idx + 8:].strip()
            lower = text.lower()
    
    # Detect PowerShell commands
    ps_prefixes = ["start-process", "start ", "get-", "set-", "new-", "remove-",
                   "invoke-", "mkdir ", "dir ", "copy ", "move ", "del "]
    for prefix in ps_prefixes:
        if lower.startswith(prefix):
            return {"tool": "run_powershell", "args": {"code": text}}
    
    # Detect Python code
    py_prefixes = ["import ", "from ", "print(", "open(", "os.", "with "]
    for prefix in py_prefixes:
        if lower.startswith(prefix):
            return {"tool": "run_python", "args": {"code": text}}
    
    return None


# ============================================================
# LOGGING
# ============================================================

def log_interaction(user_msg: str, tool_name: str, result: str):
    """Log interactions to file"""
    log_file = os.path.join(LOGS_DIR, f"log_{datetime.datetime.now().strftime('%Y-%m-%d')}.jsonl")
    entry = {
        "timestamp": datetime.datetime.now().isoformat(),
        "user": user_msg,
        "tool": tool_name,
        "result": result[:500]
    }
    with open(log_file, 'a', encoding='utf-8') as f:
        f.write(json.dumps(entry, ensure_ascii=False) + '\n')


# ============================================================
# MAIN LOOP
# ============================================================

def print_banner():
    print("""
 ================================================================
     Jane's AI Agent v2 - Local AI Assistant
     Model: {model}
     Tools: {tool_count} available
 ================================================================
 
 Commands:
   /tools     - List all available tools
   /history   - Show conversation history
   /clear     - Clear conversation history
   /install   - Install missing Python packages
   /help      - Show this help
   /exit      - Quit
 
 Just type what you need in any language!
 ================================================================
""".format(model=MODEL, tool_count=len(TOOLS)))


def handle_command(cmd: str, history: list) -> bool:
    """Handle special commands. Returns True if command was handled."""
    if cmd == "/exit" or cmd == "/quit":
        print("\nBye, Jane!")
        return True
    
    elif cmd == "/tools":
        print("\nAvailable tools:")
        for name in sorted(TOOLS.keys()):
            print(f"  - {name}")
        print()
    
    elif cmd == "/history":
        if not history:
            print("No history yet.\n")
        else:
            for msg in history[-10:]:
                role = msg["role"].upper()
                content = msg["content"][:200]
                print(f"  [{role}] {content}")
            print()
    
    elif cmd == "/clear":
        history.clear()
        print("History cleared.\n")
    
    elif cmd == "/install":
        print("Installing recommended packages...")
        packages = ["python-docx", "reportlab", "openpyxl", "python-pptx", 
                     "psutil", "playwright", "requests"]
        for pkg in packages:
            print(f"  Installing {pkg}...")
            subprocess.run([sys.executable, "-m", "pip", "install", pkg], 
                         capture_output=True)
        print("  Installing Playwright browsers...")
        subprocess.run([sys.executable, "-m", "playwright", "install", "chromium"],
                      capture_output=True)
        print("Done! All packages installed.\n")
    
    elif cmd == "/help":
        print_banner()
    
    else:
        return False  # Not a command
    
    return True  # Command was handled (but don't exit)


def main():
    print_banner()
    
    # Check Ollama connection
    try:
        resp = requests.get(f"{OLLAMA_URL}/api/tags", timeout=5)
        models = [m["name"] for m in resp.json().get("models", [])]
        if MODEL not in models and f"{MODEL}:latest" not in models:
            print(f"  WARNING: Model '{MODEL}' not found in Ollama.")
            print(f"  Available models: {', '.join(models)}")
            print(f"  Run: ollama pull {MODEL}\n")
        else:
            print(f"  Ollama connected. Model '{MODEL}' ready.\n")
    except:
        print("  WARNING: Cannot connect to Ollama. Is it running?")
        print("  Start Ollama first, then restart this agent.\n")
    
    history = []
    
    while True:
        try:
            user_input = input("You: ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nBye, Jane!")
            break
        
        if not user_input:
            continue
        
        # Handle special commands
        if user_input.startswith("/"):
            result = handle_command(user_input, history)
            if user_input in ["/exit", "/quit"]:
                break
            continue
        
        # Ask LLM
        print("  thinking...", end="\r")
        answer = ask_ollama(user_input, history)
        print("             ", end="\r")  # Clear "thinking..."
        
        # Parse response
        parsed = extract_json(answer)
        
        if not parsed:
            parsed = guess_tool(answer)
        
        if parsed and "tool" in parsed:
            tool_name = parsed["tool"]
            tool_args = parsed.get("args", {})
            
            if tool_name == "chat":
                message = tool_args.get("message", answer)
                print(f"\n  Assistant: {message}\n")
                log_interaction(user_input, "chat", message)
            
            elif tool_name in TOOLS:
                print(f"\n  [{tool_name}]", end="")
                if tool_args:
                    # Show what's being done (abbreviated)
                    brief = str(tool_args)[:100]
                    print(f" {brief}")
                else:
                    print()
                
                result = TOOLS[tool_name](tool_args)
                print(f"\n  Result:\n{result}\n")
                
                # Feed result back to history
                history.append({"role": "user", "content": f"[Tool result for {tool_name}]: {result[:1000]}"})
                log_interaction(user_input, tool_name, result)
            
            else:
                print(f"\n  Unknown tool: {tool_name}")
                print(f"  Raw response: {answer[:300]}\n")
        
        else:
            # No tool detected, show raw response
            print(f"\n  Assistant: {answer}\n")
            log_interaction(user_input, "raw", answer)


if __name__ == "__main__":
    main()
