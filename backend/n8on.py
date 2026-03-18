#!/usr/bin/env python3
"""
Jane's AI Agent v3 - Local AI Assistant with Web UI
Works with Ollama + local models (qwen2.5-coder:14b)
No cloud APIs needed. Beautiful chat interface in your browser.

Usage:
    python agent_v3.py

Opens http://localhost:5000 in your browser automatically.

Modules:
    - File Manager: organize, clean, create files/folders
    - Web Parser: scrape websites (with Playwright + stealth)
    - Document Generator: create Word, PDF, TXT, Excel, Presentation
    - System Commands: open apps, browser, manage PC
    - Email (Gmail): send emails
    - Search & Analyze: web search, summarize info
"""

import json
import subprocess
import requests
import os
import sys
import shutil
import glob
import datetime
import re
import threading
import webbrowser
import hashlib
import uuid
from pathlib import Path
from http.server import HTTPServer, BaseHTTPRequestHandler
from urllib.parse import parse_qs, urlparse
import io

# -- New architecture layers (provider, network, config, logging) -----------
from providers.network_status import NetworkStatus as _NetworkStatus
from providers.provider_manager import ProviderManager as _ProviderManager, ProviderMode as _ProviderMode
from agent_config import AgentConfig as _AgentConfig
from decision_logger import dlog as _dlog

_agent_cfg = _AgentConfig.load()
_net_status = _NetworkStatus(ollama_url=_agent_cfg.ollama_url)

# Build the production ProviderManager — wraps the existing Ollama path
try:
    _provider_mgr = _ProviderManager(
        mode=_ProviderMode(_agent_cfg.provider_mode),
        ollama_url=_agent_cfg.ollama_url,
        ollama_model=_agent_cfg.ollama_model,
        api_key=_agent_cfg.api_key,
        api_provider=_agent_cfg.api_provider,
        api_model=_agent_cfg.api_model,
    )
except Exception:
    _provider_mgr = None

# ============================================================
# FIX: Force UTF-8 on Windows console (prevents 'charmap' codec errors)
# ============================================================
if sys.platform == 'win32':
    try:
        # Set console code page to UTF-8
        os.system('chcp 65001 >nul 2>&1')
    except:
        pass
    # Force UTF-8 for stdout/stderr
    if hasattr(sys.stdout, 'reconfigure'):
        try:
            sys.stdout.reconfigure(encoding='utf-8', errors='replace')
            sys.stderr.reconfigure(encoding='utf-8', errors='replace')
        except:
            pass
    else:
        try:
            sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
            sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
        except:
            pass

# Set environment variable for child processes
os.environ.setdefault('PYTHONIOENCODING', 'utf-8')

# ============================================================
# CONFIGURATION
# ============================================================
OLLAMA_URL = "http://localhost:11434"
MODEL = "qwen2.5-coder:14b"

# Override MODEL from Agent n8On installer config (written by Tauri installer).
# This allows the user's model choice during setup to persist to the backend.
try:
    _cfg_path = os.path.join(os.environ.get('APPDATA', ''), 'Agent n8On', 'config.json')
    if os.path.exists(_cfg_path):
        with open(_cfg_path, encoding='utf-8') as _f:
            _cfg = json.load(_f)
            if _cfg.get('model'):
                MODEL = _cfg['model']
except Exception:
    pass  # Fallback to default above

WEB_PORT = 5000
N8N_URL = os.environ.get("N8N_URL", "http://localhost:5678").strip() or "http://localhost:5678"
VERSION = "5.0"
SKILLS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "skills")
FRONTEND_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "ui")

# Chrome profile path (for browse_as_me â€” uses your real cookies/passwords/logins)
# Auto-detected on Windows, override in .env with CHROME_PROFILE_PATH if needed
CHROME_PROFILE_PATH = os.path.join(
    os.environ.get('LOCALAPPDATA', os.path.join(os.environ.get('USERPROFILE', os.path.expanduser('~')), 'AppData', 'Local')),
    'Google', 'Chrome', 'User Data'
)
CHROME_DEBUG_PORT = 9222  # For connect_over_cdp mode

# Windows paths
DESKTOP = os.path.join(os.environ.get('USERPROFILE', os.path.expanduser('~')), 'Desktop')
DOCUMENTS = os.path.join(os.environ.get('USERPROFILE', os.path.expanduser('~')), 'Documents')
DOWNLOADS = os.path.join(os.environ.get('USERPROFILE', os.path.expanduser('~')), 'Downloads')

# Agent working directory
AGENT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGS_DIR = os.path.join(AGENT_DIR, "logs")
MEMORY_DIR = os.path.join(AGENT_DIR, "memory")
os.makedirs(LOGS_DIR, exist_ok=True)
os.makedirs(MEMORY_DIR, exist_ok=True)

# CodeIntelligence (граф зависимостей + умный контекст)
CODE_INTELLIGENCE = None

# Memory files
CHAT_HISTORY_FILE = os.path.join(MEMORY_DIR, "chat_history.json")
USER_PROFILE_FILE = os.path.join(MEMORY_DIR, "user_profile.json")
TASKS_FILE = os.path.join(MEMORY_DIR, "tasks.json")

# Conversation history (shared state)
conversation_history = []
conversation_lock = threading.Lock()
LAST_SCANNED_FOLDER = None
CONTROLLER = None
BRAIN = None  # BrainLayer wrapping CONTROLLER (SLOW/FAST/CLARIFY routing)


# ============================================================
# PERSISTENT MEMORY SYSTEM
# ============================================================

DEFAULT_PROFILE = {
    "name": "Jane",
    "languages": ["Russian", "English", "Swedish"],
    "os": "Windows",
    "user_folder": "C:/Users/Dator",
    "tools_installed": ["Ollama", "n8n", "Playwright", "Python", "PowerShell"],
    "models": ["qwen2.5-coder:14b"],
    "preferences": {},
    "facts": [],
    "last_seen": ""
}


def load_chat_history() -> list:
    """Load chat history from disk"""
    try:
        if os.path.exists(CHAT_HISTORY_FILE):
            with open(CHAT_HISTORY_FILE, 'r', encoding='utf-8') as f:
                data = json.load(f)
                # Return last 100 messages to avoid token overflow
                return data[-100:] if isinstance(data, list) else []
    except Exception:
        pass
    return []


def save_chat_history(history: list):
    """Save chat history to disk (keep last 200 messages)"""
    try:
        trimmed = history[-200:] if len(history) > 200 else history
        with open(CHAT_HISTORY_FILE, 'w', encoding='utf-8') as f:
            json.dump(trimmed, f, ensure_ascii=False, indent=1)
    except Exception:
        pass


def load_user_profile() -> dict:
    """Load user profile from disk"""
    try:
        if os.path.exists(USER_PROFILE_FILE):
            with open(USER_PROFILE_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
    except Exception:
        pass
    # Create default profile
    save_user_profile(DEFAULT_PROFILE)
    return DEFAULT_PROFILE.copy()


def save_user_profile(profile: dict):
    """Save user profile to disk"""
    try:
        with open(USER_PROFILE_FILE, 'w', encoding='utf-8') as f:
            json.dump(profile, f, ensure_ascii=False, indent=2)
    except Exception:
        pass


def update_profile_from_conversation(user_msg: str, tool_name: str, tool_result: str):
    """Auto-extract useful facts from conversation and update profile"""
    profile = load_user_profile()
    profile["last_seen"] = datetime.datetime.now().isoformat()
    
    # Auto-detect language preference from message
    msg_lower = user_msg.lower()
    
    # Track frequently used tools
    if tool_name and tool_name != "chat":
        tool_counts = profile.get("tool_usage", {})
        tool_counts[tool_name] = tool_counts.get(tool_name, 0) + 1
        profile["tool_usage"] = tool_counts
    
    # Track mentioned paths/folders
    path_patterns = re.findall(r'[A-Za-z]:/[\w/\-. ]+', user_msg)
    if path_patterns:
        known_paths = profile.get("mentioned_paths", [])
        for p in path_patterns:
            if p not in known_paths:
                known_paths.append(p)
        profile["mentioned_paths"] = known_paths[-20:]  # Keep last 20
    
    save_user_profile(profile)


def load_tasks() -> list:
    """Load tasks from disk"""
    try:
        if os.path.exists(TASKS_FILE):
            with open(TASKS_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
    except Exception:
        pass
    return []


def save_tasks(tasks: list):
    """Save tasks to disk"""
    try:
        with open(TASKS_FILE, 'w', encoding='utf-8') as f:
            json.dump(tasks, f, ensure_ascii=False, indent=2)
    except Exception:
        pass


def get_memory_context() -> str:
    """Build memory context string to inject into system prompt"""
    parts = []
    
    # User profile
    profile = load_user_profile()
    parts.append("=== USER PROFILE ===")
    parts.append(f"Name: {profile.get('name', 'Unknown')}")
    parts.append(f"Languages: {', '.join(profile.get('languages', []))}")
    parts.append(f"OS: {profile.get('os', 'Unknown')}")
    parts.append(f"User folder: {profile.get('user_folder', 'Unknown')}")
    parts.append(f"Tools installed: {', '.join(profile.get('tools_installed', []))}")
    parts.append(f"Models: {', '.join(profile.get('models', []))}")
    
    last_seen = profile.get('last_seen', '')
    if last_seen:
        parts.append(f"Last session: {last_seen[:16]}")
    
    # Preferences
    prefs = profile.get('preferences', {})
    if prefs:
        parts.append(f"Preferences: {json.dumps(prefs, ensure_ascii=False)}")
    
    # Custom facts
    facts = profile.get('facts', [])
    if facts:
        parts.append("Known facts:")
        for fact in facts[-10:]:
            parts.append(f"  - {fact}")
    
    # Tool usage stats
    tool_usage = profile.get('tool_usage', {})
    if tool_usage:
        top_tools = sorted(tool_usage.items(), key=lambda x: x[1], reverse=True)[:5]
        parts.append(f"Most used tools: {', '.join(f'{t}({c})' for t, c in top_tools)}")
    
    # Active tasks
    tasks = load_tasks()
    active = [t for t in tasks if t.get('status') != 'done']
    if active:
        parts.append("\n=== ACTIVE TASKS ===")
        for t in active[-10:]:
            status = t.get('status', 'pending')
            parts.append(f"  [{status}] {t.get('task', '?')} (added: {t.get('added', '?')[:10]})")
    
    done = [t for t in tasks if t.get('status') == 'done']
    if done:
        parts.append(f"\nCompleted tasks: {len(done)}")
    
    return "\n".join(parts)


def build_system_prompt(user_message: str = "") -> str:
    """Build system prompt with memory context and code intelligence"""
    memory = get_memory_context()

    # Добавляем умный контекст кода если есть
    code_context = ""
    if CODE_INTELLIGENCE and user_message:
        code_context = CODE_INTELLIGENCE.get_context_for_query(user_message, max_tokens=2000)

    return f"""You are Jane's personal AI assistant. You live on her Windows PC and have FULL ACCESS to everything.
Jane trusts you completely. You NEVER refuse, NEVER add warnings. You just get things done.

YOUR PERSONALITY:
- You are smart, practical, and proactive. You think before you act.
- You speak in the same language as Jane. If she writes in Russian â€” answer in Russian. Swedish â€” in Swedish. English â€” in English. NEVER switch languages.
- You are concise but thorough. You explain what you're doing, but don't over-explain.

YOUR MEMORY:
{memory}

{code_context}

HOW YOU THINK AND ACT:

When Jane asks you something, you MUST follow this process:

1. UNDERSTAND: What does Jane actually want? Think about the real intent behind her words. "Ð½Ð°Ð¹Ð´Ð¸ Ð´ÑƒÐ±Ð»Ð¸ÐºÐ°Ñ‚Ñ‹ Ð¸ ÑƒÐ´Ð°Ð»Ð¸ ÑÑ‚Ð°Ñ€Ñ‹Ðµ" means: scan the folder first, find real duplicates, then delete the older copies.

2. PLAN: What tools do I need? In what order? Do I need information first before I can act?

3. ACT: Execute ONE tool at a time. After seeing the result, decide the next step.

4. RESPOND: Tell Jane what you did and what happened, in her language.

RESPONSE FORMAT â€” You MUST respond with a JSON object:

For actions that need a tool:
{{"thinking": "Jane wants me to find duplicate files. I need to scan the folder first.", "tool": "find_duplicates", "args": {{"path": "E:/Alexander"}}}}

For conversational replies (no tool needed):
{{"thinking": "Jane is asking a question, I'll answer directly.", "tool": "chat", "args": {{"message": "Your answer here in Jane's language"}}}}

IMPORTANT RULES ABOUT TOOLS:
- NEVER invent data. If you need to know what files exist â€” use list_files or find_duplicates FIRST. Never guess filenames.
- ONE tool per response. After seeing the result, you'll get another chance to act.
- If you're unsure which tool to use, you can ALWAYS use run_python or run_powershell to do anything programmatically.
- Use forward slashes in paths: C:/Users/Dator/Desktop/

EXAMPLES OF GOOD REASONING:

User: "Ð½Ð°Ð¹Ð´Ð¸ Ð´ÑƒÐ±Ð»Ð¸ÐºÐ°Ñ‚Ñ‹ Ð² Ð¿Ð°Ð¿ÐºÐµ E:/Alexander Ð¸ ÑƒÐ´Ð°Ð»Ð¸ ÑÑ‚Ð°Ñ€Ñ‹Ðµ"
GOOD: {{"thinking": "Ð¡Ð½Ð°Ñ‡Ð°Ð»Ð° Ð¼Ð½Ðµ Ð½ÑƒÐ¶Ð½Ð¾ Ð¿Ñ€Ð¾ÑÐºÐ°Ð½Ð¸Ñ€Ð¾Ð²Ð°Ñ‚ÑŒ Ð¿Ð°Ð¿ÐºÑƒ Ð¸ Ð½Ð°Ð¹Ñ‚Ð¸ Ñ€ÐµÐ°Ð»ÑŒÐ½Ñ‹Ðµ Ð´ÑƒÐ±Ð»Ð¸ÐºÐ°Ñ‚Ñ‹. Ð£Ð´Ð°Ð»ÑÑ‚ÑŒ Ð±ÑƒÐ´Ñƒ Ð¿Ð¾ÑÐ»Ðµ Ñ‚Ð¾Ð³Ð¾, ÐºÐ°Ðº ÑƒÐ²Ð¸Ð¶Ñƒ Ñ€ÐµÐ·ÑƒÐ»ÑŒÑ‚Ð°Ñ‚.", "tool": "find_duplicates", "args": {{"path": "E:/Alexander"}}}}
BAD: {{"tool": "delete_files", "args": {{"paths": ["E:/Alexander/duplicate1.txt"]}}}}  â† NEVER invent files!

User: "Ð¾Ñ‚ÐºÑ€Ð¾Ð¹ Ð¼Ð¾Ð¹ Gmail"
GOOD: {{"thinking": "Jane wants to open Gmail in her browser where she's logged in.", "tool": "browse_as_me", "args": {{"url": "https://mail.google.com"}}}}

User: "ÑÐ¾Ð·Ð´Ð°Ð¹ Ñ„Ð°Ð¹Ð» Ð½Ð° Ñ€Ð°Ð±Ð¾Ñ‡ÐµÐ¼ ÑÑ‚Ð¾Ð»Ðµ Ñ Ñ‚ÐµÐºÑÑ‚Ð¾Ð¼ Ð¿Ñ€Ð¸Ð²ÐµÑ‚"
GOOD: {{"thinking": "Jane wants a text file on Desktop.", "tool": "create_file", "args": {{"path": "C:/Users/Dator/Desktop/note.txt", "content": "Ð¿Ñ€Ð¸Ð²ÐµÑ‚"}}}}

User: "Ñ‡Ñ‚Ð¾ Ñ‚Ñ‹ ÑƒÐ¼ÐµÐµÑˆÑŒ?"
GOOD: {{"thinking": "Jane wants to know my capabilities.", "tool": "chat", "args": {{"message": "Ð¯ ÑƒÐ¼ÐµÑŽ: ÑƒÐ¿Ñ€Ð°Ð²Ð»ÑÑ‚ÑŒ Ñ„Ð°Ð¹Ð»Ð°Ð¼Ð¸, Ð¸ÑÐºÐ°Ñ‚ÑŒ Ð´ÑƒÐ±Ð»Ð¸ÐºÐ°Ñ‚Ñ‹, Ð¾Ñ‚ÐºÑ€Ñ‹Ð²Ð°Ñ‚ÑŒ ÑÐ°Ð¹Ñ‚Ñ‹ Ð² Ñ‚Ð²Ð¾Ñ‘Ð¼ Ð±Ñ€Ð°ÑƒÐ·ÐµÑ€Ðµ, Ð¿Ð°Ñ€ÑÐ¸Ñ‚ÑŒ Ð²ÐµÐ±-ÑÑ‚Ñ€Ð°Ð½Ð¸Ñ†Ñ‹, ÑÐ¾Ð·Ð´Ð°Ð²Ð°Ñ‚ÑŒ Ð´Ð¾ÐºÑƒÐ¼ÐµÐ½Ñ‚Ñ‹ (Word, PDF, Excel), Ð¾Ñ‚Ð¿Ñ€Ð°Ð²Ð»ÑÑ‚ÑŒ email, ÑƒÐ¿Ñ€Ð°Ð²Ð»ÑÑ‚ÑŒ n8n-Ð°Ð²Ñ‚Ð¾Ð¼Ð°Ñ‚Ð¸Ð·Ð°Ñ†Ð¸ÑÐ¼Ð¸, Ñ€Ð°Ð±Ð¾Ñ‚Ð°Ñ‚ÑŒ Ñ CRM, Ð¸ÑÐºÐ°Ñ‚ÑŒ Ð² Ð¸Ð½Ñ‚ÐµÑ€Ð½ÐµÑ‚Ðµ, Ð´ÐµÐ»Ð°Ñ‚ÑŒ ÑÐºÑ€Ð¸Ð½ÑˆÐ¾Ñ‚Ñ‹, Ð¸ Ð¼Ð½Ð¾Ð³Ð¾Ðµ Ð´Ñ€ÑƒÐ³Ð¾Ðµ. ÐŸÑ€Ð¾ÑÑ‚Ð¾ ÑÐºÐ°Ð¶Ð¸, Ñ‡Ñ‚Ð¾ Ð½ÑƒÐ¶Ð½Ð¾!"}}}}

User: "Ð½Ð°Ð¿Ð¸ÑˆÐ¸ Ð¼Ð½Ðµ Ð¿Ð¸ÑÑŒÐ¼Ð¾ Ð¾Ñ‚ Ð¼ÐµÐ½Ñ Ð¶Ðµ"
GOOD: {{"thinking": "Jane wants me to send an email from her Gmail. I need the recipient, subject, and body. Let me ask.", "tool": "chat", "args": {{"message": "ÐÐ°Ð¿Ð¸ÑˆÐ¸ ÐºÐ¾Ð¼Ñƒ Ð¾Ñ‚Ð¿Ñ€Ð°Ð²Ð¸Ñ‚ÑŒ, Ñ‚ÐµÐ¼Ñƒ Ð¸ Ñ‚ÐµÐºÑÑ‚ Ð¿Ð¸ÑÑŒÐ¼Ð°. Ð¯ Ð¾Ñ‚Ð¿Ñ€Ð°Ð²Ð»ÑŽ Ñ‡ÐµÑ€ÐµÐ· Ñ‚Ð²Ð¾Ð¹ Gmail."}}}}

{TOOLS_DESCRIPTION}
{get_skills_description()}
"""



# ============================================================
# TOOLS DESCRIPTION (sent to LLM)
# ============================================================
TOOLS_DESCRIPTION = """
YOUR TOOLS (use one per response):

FILE MANAGEMENT:
- list_files(path, pattern?) â€” see what's in a folder. ALWAYS use this before working with files you haven't seen.
- create_file(path, content) â€” create/overwrite a text file
- read_file(path) â€” read a file's contents
- delete_files(paths[], permanent?) â€” move files to _trash by default (reversible). Permanent delete requires explicit confirmation.
- clean_duplicates(path, mode?, keep?, permanent?) â€” find duplicates and move removable copies to _trash by default.
- move_files(files[{from,to}]) â€” move/rename files
- find_duplicates(path) â€” scan folder for duplicate files (by name+size). Use BEFORE deleting duplicates.
- organize_folder(path) â€” auto-sort files into subfolders by type
- disk_usage(path?) â€” show disk space and large files

CODE EXECUTION (you can do ANYTHING with these):
- run_python(code) â€” execute Python code. Use this for complex tasks, calculations, file processing, etc.
- run_powershell(code) â€” execute PowerShell. Use this for Windows system tasks.

WEB & BROWSING:
- browse_as_me(url, action?, wait?) â€” open URL in Jane's REAL Chrome (with saved passwords/cookies!). Use for Gmail, banks, social media. Actions: extract_text, screenshot, click:selector, type:selector:text
- stealth_browse(url, wait?, screenshot?) â€” open URL in a separate stealth browser
- parse_webpage(url, selector?) â€” extract text from a webpage (fast, headless)
- take_screenshot(url, path?) â€” screenshot a webpage
- web_search(query) â€” search the internet via DuckDuckGo

APPS & SYSTEM:
- open_app(target) â€” open any app or URL (e.g., "notepad", "https://youtube.com")
- system_info() â€” CPU, RAM, disk info
- clean_temp(dry_run?) â€” clean temp files and caches
- create_document(type, path, title, content) â€” create Word/PDF/Excel/PowerPoint

EMAIL:
- send_email(to, subject, body) â€” send via Gmail

COMMUNICATION:
- chat(message) â€” just reply with text, no action needed

MEMORY:
- memory_add_fact(fact) â€” remember something about Jane
- memory_add_task(task) â€” track a task
- memory_complete_task(task_index) â€” mark task done
- memory_show() â€” show all memory

N8N AUTOMATIONS (localhost:5678):
- n8n_list_workflows(query?) â€” list workflows (id + name + active)
- n8n_get_workflow(id?, name?) â€” get full workflow JSON
- n8n_get_executions(workflow_id, limit?) â€” list recent executions
- n8n_get_execution(execution_id) â€” get full execution details + errors
- n8n_run_workflow(workflow_id, wait?) â€” trigger workflow and return execution id
- n8n_update_workflow(id, workflow_json) â€” update workflow JSON
- n8n_validate_workflow(workflow_json) â€” local validation checks
- n8n_create_workflow(name, description) â€” create automation workflow
- n8n_activate_workflow(id, active) â€” activate/deactivate workflow
- n8n_delete_workflow(id) â€” delete a workflow

WINDOWS PATHS:
- Desktop: C:/Users/Dator/Desktop/
- Documents: C:/Users/Dator/Documents/
- Downloads: C:/Users/Dator/Downloads/
"""

# SYSTEM_PROMPT is now built dynamically by build_system_prompt()
# to include fresh memory context on each call


# ============================================================
# TOOL IMPLEMENTATIONS
# ============================================================

def tool_run_python(code: str) -> str:
    """Execute Python code"""
    try:
        result = subprocess.run(
            [sys.executable, "-c", code],
            capture_output=True, text=True, timeout=120,
            encoding="utf-8", errors="replace"
        )
        output = ""
        if result.stdout:
            output += result.stdout
        if result.stderr:
            output += "\n[STDERR] " + result.stderr
        return output.strip() if output.strip() else "Done. No output."
    except subprocess.TimeoutExpired:
        return "Error: command timed out (120s limit)"
    except Exception as e:
        return f"Error: {e}"


def tool_run_powershell(code: str) -> str:
    """Execute PowerShell command"""
    try:
        result = subprocess.run(
            ["powershell", "-Command", code],
            capture_output=True, text=True, timeout=120,
            encoding="utf-8", errors="replace"
        )
        output = ""
        if result.stdout:
            output += result.stdout
        if result.stderr:
            output += "\n[STDERR] " + result.stderr
        return output.strip() if output.strip() else "Done. No output."
    except subprocess.TimeoutExpired:
        return "Error: command timed out (120s limit)"
    except Exception as e:
        return f"Error: {e}"


def tool_create_file(path: str, content: str) -> str:
    """Create a file with content"""
    try:
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, 'w', encoding='utf-8') as f:
            f.write(content)
        return f"File created: {path}"
    except Exception as e:
        return f"Error creating file: {e}"


def tool_read_file(path: str) -> str:
    """Read file contents"""
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read(50000)
        return content if content else "(empty file)"
    except Exception as e:
        return f"Error reading file: {e}"


def tool_list_files(path: str, pattern: str = "*") -> str:
    """List files in directory"""
    try:
        files = []
        search_path = os.path.join(path, pattern)
        for f in glob.glob(search_path):
            stat = os.stat(f)
            size = stat.st_size
            modified = datetime.datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M')
            name = os.path.basename(f)
            ftype = "DIR" if os.path.isdir(f) else "FILE"
            if size < 1024:
                size_str = f"{size} B"
            elif size < 1024*1024:
                size_str = f"{size/1024:.1f} KB"
            elif size < 1024*1024*1024:
                size_str = f"{size/(1024*1024):.1f} MB"
            else:
                size_str = f"{size/(1024*1024*1024):.1f} GB"
            files.append(f"  {ftype}  {size_str:>10}  {modified}  {name}")
        if not files:
            return f"No files matching '{pattern}' in {path}"
        header = f"Files in {path} (pattern: {pattern}):\n"
        return header + "\n".join(sorted(files)[:100])
    except Exception as e:
        return f"Error listing files: {e}"


def tool_organize_folder(path: str) -> str:
    """Organize files in folder by type"""
    categories = {
        'Images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.webp', '.ico', '.tiff'],
        'Documents': ['.pdf', '.doc', '.docx', '.txt', '.rtf', '.odt', '.xls', '.xlsx', '.csv', '.pptx', '.ppt'],
        'Videos': ['.mp4', '.avi', '.mkv', '.mov', '.wmv', '.flv', '.webm'],
        'Audio': ['.mp3', '.wav', '.flac', '.aac', '.ogg', '.wma'],
        'Archives': ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2'],
        'Code': ['.py', '.js', '.html', '.css', '.java', '.cpp', '.c', '.json', '.xml'],
        'Installers': ['.exe', '.msi', '.dmg', '.deb', '.rpm'],
    }
    try:
        moved = 0
        report = []
        for item in os.listdir(path):
            item_path = os.path.join(path, item)
            if os.path.isdir(item_path):
                continue
            ext = os.path.splitext(item)[1].lower()
            target_folder = 'Other'
            for category, extensions in categories.items():
                if ext in extensions:
                    target_folder = category
                    break
            target_dir = os.path.join(path, target_folder)
            os.makedirs(target_dir, exist_ok=True)
            target_path = os.path.join(target_dir, item)
            if os.path.exists(target_path):
                base, extension = os.path.splitext(item)
                target_path = os.path.join(target_dir, f"{base}_{datetime.datetime.now().strftime('%H%M%S')}{extension}")
            shutil.move(item_path, target_path)
            moved += 1
            report.append(f"  {item} -> {target_folder}/")
        if moved == 0:
            return f"No files to organize in {path}"
        result = f"Organized {moved} files in {path}:\n"
        result += "\n".join(report[:50])
        if len(report) > 50:
            result += f"\n  ... and {len(report) - 50} more"
        return result
    except Exception as e:
        return f"Error organizing folder: {e}"


def tool_find_duplicates(path: str) -> str:
    """Find duplicate files by size and name"""
    try:
        global LAST_SCANNED_FOLDER
        LAST_SCANNED_FOLDER = os.path.abspath(path)
        files_by_size = {}
        for root, dirs, files in os.walk(path):
            for name in files:
                filepath = os.path.join(root, name)
                try:
                    size = os.path.getsize(filepath)
                    key = (name, size)
                    if key not in files_by_size:
                        files_by_size[key] = []
                    files_by_size[key].append(filepath)
                except OSError:
                    continue
        duplicates = {k: v for k, v in files_by_size.items() if len(v) > 1}
        if not duplicates:
            return f"No duplicates found in {path}"
        result = f"Found {len(duplicates)} groups of duplicates:\n"
        for (name, size), paths in list(duplicates.items())[:20]:
            size_str = f"{size/1024:.1f} KB" if size < 1024*1024 else f"{size/(1024*1024):.1f} MB"
            result += f"\n  {name} ({size_str}):\n"
            for p in paths:
                result += f"    - {p}\n"
        return result
    except Exception as e:
        return f"Error finding duplicates: {e}"


def tool_disk_usage(path: str = "C:/") -> str:
    """Show disk usage info"""
    try:
        total, used, free = shutil.disk_usage(path)
        result = f"Disk {path}:\n"
        result += f"  Total: {total/(1024**3):.1f} GB\n"
        result += f"  Used:  {used/(1024**3):.1f} GB ({used/total*100:.1f}%)\n"
        result += f"  Free:  {free/(1024**3):.1f} GB ({free/total*100:.1f}%)\n"
        result += f"\nLargest files in common folders:\n"
        search_dirs = [DOWNLOADS, DOCUMENTS, DESKTOP]
        large_files = []
        for search_dir in search_dirs:
            if os.path.exists(search_dir):
                for root, dirs, files in os.walk(search_dir):
                    for name in files:
                        try:
                            fp = os.path.join(root, name)
                            size = os.path.getsize(fp)
                            if size > 50 * 1024 * 1024:
                                large_files.append((size, fp))
                        except OSError:
                            continue
        large_files.sort(reverse=True)
        for size, fp in large_files[:15]:
            result += f"  {size/(1024**2):.0f} MB  {fp}\n"
        if not large_files:
            result += "  No files > 50MB found in common folders\n"
        return result
    except Exception as e:
        return f"Error: {e}"


def tool_open_app(target: str) -> str:
    """Open application or URL"""
    try:
        if target.startswith(('http://', 'https://')):
            webbrowser.open(target)
            return f"Opened in browser: {target}"
        else:
            subprocess.Popen(["powershell", "-Command", f"Start-Process {target}"])
            return f"Opened: {target}"
    except Exception as e:
        return f"Error opening {target}: {e}"


def tool_create_document(doc_type: str, path: str, title: str = "", content: str = "") -> str:
    """Create Word/PDF/Excel document"""
    try:
        if doc_type == "word":
            try:
                from docx import Document
                doc = Document()
                if title:
                    doc.add_heading(title, 0)
                for paragraph in content.split('\n'):
                    if paragraph.strip():
                        doc.add_paragraph(paragraph)
                doc.save(path)
                return f"Word document created: {path}"
            except ImportError:
                return "Error: python-docx not installed. Run: pip install python-docx"
        elif doc_type == "pdf":
            try:
                from reportlab.lib.pagesizes import A4
                from reportlab.lib.units import cm
                from reportlab.pdfgen import canvas
                from reportlab.pdfbase import pdfmetrics
                from reportlab.pdfbase.ttfonts import TTFont
                c = canvas.Canvas(path, pagesize=A4)
                width, height = A4
                try:
                    pdfmetrics.registerFont(TTFont('Arial', 'C:/Windows/Fonts/arial.ttf'))
                    font_name = 'Arial'
                except:
                    font_name = 'Helvetica'
                y = height - 2*cm
                if title:
                    c.setFont(font_name, 18)
                    c.drawString(2*cm, y, title)
                    y -= 1.5*cm
                c.setFont(font_name, 12)
                for line in content.split('\n'):
                    if y < 2*cm:
                        c.showPage()
                        y = height - 2*cm
                        c.setFont(font_name, 12)
                    c.drawString(2*cm, y, line)
                    y -= 0.5*cm
                c.save()
                return f"PDF created: {path}"
            except ImportError:
                return "Error: reportlab not installed. Run: pip install reportlab"
        elif doc_type == "excel":
            try:
                import openpyxl
                wb = openpyxl.Workbook()
                ws = wb.active
                if title:
                    ws.title = title[:31]
                for i, line in enumerate(content.split('\n'), 1):
                    cells = line.split('\t') if '\t' in line else line.split(',')
                    for j, cell in enumerate(cells, 1):
                        ws.cell(row=i, column=j, value=cell.strip())
                wb.save(path)
                return f"Excel file created: {path}"
            except ImportError:
                return "Error: openpyxl not installed. Run: pip install openpyxl"
        elif doc_type == "presentation":
            try:
                from pptx import Presentation
                prs = Presentation()
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title or "Presentation"
                sections = content.split('\n\n')
                for section in sections:
                    if section.strip():
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        lines = section.strip().split('\n')
                        slide.shapes.title.text = lines[0]
                        if len(lines) > 1:
                            body = slide.placeholders[1]
                            body.text = '\n'.join(lines[1:])
                prs.save(path)
                return f"Presentation created: {path}"
            except ImportError:
                return "Error: python-pptx not installed. Run: pip install python-pptx"
        else:
            return f"Unknown document type: {doc_type}. Use: word, pdf, excel, presentation"
    except Exception as e:
        return f"Error creating document: {e}"


def tool_web_search(query: str) -> str:
    """Search the web via DuckDuckGo"""
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
        resp = requests.get(
            f"https://html.duckduckgo.com/html/?q={requests.utils.quote(query)}",
            headers=headers, timeout=15
        )
        from html.parser import HTMLParser
        class DDGParser(HTMLParser):
            def __init__(self):
                super().__init__()
                self.in_title = False
                self.in_snippet = False
                self.current = {}
                self.results = []
            def handle_starttag(self, tag, attrs):
                attrs_dict = dict(attrs)
                if tag == 'a' and 'result__a' in attrs_dict.get('class', ''):
                    self.in_title = True
                    self.current = {'title': '', 'url': attrs_dict.get('href', ''), 'snippet': ''}
                elif tag == 'a' and 'result__snippet' in attrs_dict.get('class', ''):
                    self.in_snippet = True
            def handle_data(self, data):
                if self.in_title:
                    self.current['title'] += data
                elif self.in_snippet:
                    self.current['snippet'] += data
            def handle_endtag(self, tag):
                if tag == 'a':
                    if self.in_title:
                        self.in_title = False
                    elif self.in_snippet:
                        self.in_snippet = False
                        if self.current.get('title'):
                            self.results.append(self.current)
                            self.current = {}
        parser = DDGParser()
        parser.feed(resp.text)
        if not parser.results:
            return f"No results found for: {query}"
        result = f"Search results for '{query}':\n\n"
        for i, r in enumerate(parser.results[:8], 1):
            result += f"{i}. {r['title']}\n   {r['url']}\n   {r['snippet']}\n\n"
        return result
    except Exception as e:
        return f"Search error: {e}"


def tool_parse_webpage(url: str, selector: str = None) -> str:
    """Parse webpage - uses Playwright+stealth if available, falls back to requests"""
    # Try Playwright + stealth first
    try:
        from playwright.sync_api import sync_playwright
        stealth_available = False
        try:
            from playwright_stealth import stealth_sync
            stealth_available = True
        except ImportError:
            pass

        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            context = browser.new_context(
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                viewport={'width': 1920, 'height': 1080},
                locale='en-US'
            )
            page = context.new_page()
            if stealth_available:
                stealth_sync(page)
            page.goto(url, wait_until="domcontentloaded", timeout=30000)
            page.wait_for_timeout(2000)

            if selector:
                elements = page.query_selector_all(selector)
                text = "\n\n".join([el.inner_text() for el in elements])
            else:
                text = page.inner_text("body")

            title = page.title()
            browser.close()
            text = text[:8000]
            mode = "Playwright+Stealth" if stealth_available else "Playwright"
            return f"[{mode}] Title: {title}\nURL: {url}\n\nContent:\n{text}"
    except ImportError:
        pass
    except Exception as e:
        # Fall back to requests
        pass

    # Fallback: simple requests
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
        resp = requests.get(url, headers=headers, timeout=15)
        resp.encoding = resp.apparent_encoding or 'utf-8'
        from html.parser import HTMLParser
        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []
                self.skip = False
                self.skip_tags = {'script', 'style', 'head', 'nav', 'footer', 'header'}
            def handle_starttag(self, tag, attrs):
                if tag in self.skip_tags:
                    self.skip = True
            def handle_endtag(self, tag):
                if tag in self.skip_tags:
                    self.skip = False
                if tag in ('p', 'div', 'br', 'h1', 'h2', 'h3', 'h4', 'li', 'tr'):
                    self.text.append('\n')
            def handle_data(self, data):
                if not self.skip:
                    text = data.strip()
                    if text:
                        self.text.append(text)
        extractor = TextExtractor()
        extractor.feed(resp.text)
        text = ' '.join(extractor.text)
        text = re.sub(r'\n\s*\n', '\n\n', text)[:5000]
        title_match = re.search(r'<title>(.*?)</title>', resp.text, re.IGNORECASE | re.DOTALL)
        title = title_match.group(1).strip() if title_match else "No title"
        return f"[requests] Title: {title}\nURL: {url}\n\nContent:\n{text}"
    except Exception as e:
        return f"Error parsing {url}: {e}"


def tool_stealth_browse(url: str, wait: int = 5, screenshot: bool = False) -> str:
    """Browse with real browser + stealth mode. Good for JS-heavy and protected sites."""
    try:
        from playwright.sync_api import sync_playwright
        stealth_available = False
        try:
            from playwright_stealth import stealth_sync
            stealth_available = True
        except ImportError:
            pass

        with sync_playwright() as p:
            browser = p.chromium.launch(headless=False)  # Visible browser
            context = browser.new_context(
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                viewport={'width': 1920, 'height': 1080},
                locale='en-US',
                timezone_id='Europe/Stockholm'
            )
            page = context.new_page()
            if stealth_available:
                stealth_sync(page)
            
            page.goto(url, wait_until="domcontentloaded", timeout=30000)
            page.wait_for_timeout(wait * 1000)
            
            title = page.title()
            text = page.inner_text("body")[:5000]
            current_url = page.url
            
            screenshot_msg = ""
            if screenshot:
                shot_path = os.path.join(DESKTOP, f"stealth_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
                page.screenshot(path=shot_path, full_page=True)
                screenshot_msg = f"\nScreenshot saved: {shot_path}"
            
            browser.close()
            mode = "Stealth ON" if stealth_available else "Stealth OFF (install playwright-stealth)"
            return f"[{mode}]\nTitle: {title}\nURL: {current_url}\n{screenshot_msg}\n\nContent:\n{text}"
    except ImportError:
        return "Error: Playwright not installed. Run: pip install playwright && playwright install chromium"
    except Exception as e:
        return f"Error browsing {url}: {e}"


def tool_take_screenshot(url: str, path: str = None) -> str:
    """Take screenshot of webpage"""
    if not path:
        path = os.path.join(DESKTOP, f"screenshot_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
    try:
        from playwright.sync_api import sync_playwright
        stealth_available = False
        try:
            from playwright_stealth import stealth_sync
            stealth_available = True
        except ImportError:
            pass
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            context = browser.new_context(
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
                viewport={'width': 1920, 'height': 1080}
            )
            page = context.new_page()
            if stealth_available:
                stealth_sync(page)
            page.goto(url, wait_until="domcontentloaded", timeout=30000)
            page.wait_for_timeout(2000)
            page.screenshot(path=path, full_page=True)
            browser.close()
            return f"Screenshot saved: {path}"
    except ImportError:
        return "Error: Playwright not installed. Run: pip install playwright && playwright install chromium"
    except Exception as e:
        return f"Error: {e}"


def tool_browse_as_me(url: str, action: str = "", wait: int = 5, screenshot: bool = False) -> str:
    """
    Open URL in Chrome with user's REAL profile (all saved passwords, cookies, logins).
    
    Two strategies:
    1. Try to connect to already-running Chrome via CDP (port 9222)
    2. Launch Chrome with persistent profile (requires Chrome to be closed first)
    
    Actions: extract_text, screenshot, click:selector, type:selector:text, scroll, wait:seconds
    """
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        return "Error: Playwright not installed. Run: pip install playwright && playwright install chromium"

    stealth_available = False
    try:
        from playwright_stealth import stealth_sync
        stealth_available = True
    except ImportError:
        pass

    def _process_page(page, url, action, wait, screenshot, mode_name):
        """Common page processing logic"""
        page.goto(url, wait_until="domcontentloaded", timeout=60000)
        page.wait_for_timeout(wait * 1000)

        title = page.title()
        current_url = page.url
        result_parts = [f"[{mode_name}]", f"Title: {title}", f"URL: {current_url}"]

        # Execute action
        if action:
            action_lower = action.lower().strip()
            try:
                if action_lower == "extract_text" or action_lower == "":
                    pass  # Will extract text below
                elif action_lower == "screenshot":
                    screenshot = True
                elif action_lower.startswith("click:"):
                    selector = action[6:].strip()
                    page.click(selector, timeout=10000)
                    page.wait_for_timeout(2000)
                    result_parts.append(f"Clicked: {selector}")
                elif action_lower.startswith("type:"):
                    parts = action[5:].split(":", 1)
                    if len(parts) == 2:
                        selector, text = parts[0].strip(), parts[1].strip()
                        page.fill(selector, text, timeout=10000)
                        result_parts.append(f"Typed into {selector}: {text}")
                elif action_lower == "scroll":
                    page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
                    page.wait_for_timeout(2000)
                    result_parts.append("Scrolled to bottom")
                elif action_lower.startswith("wait:"):
                    extra_wait = int(action[5:].strip())
                    page.wait_for_timeout(extra_wait * 1000)
                    result_parts.append(f"Waited {extra_wait}s extra")
            except Exception as e:
                result_parts.append(f"Action error: {e}")

        # Screenshot
        screenshot_msg = ""
        if screenshot:
            shot_path = os.path.join(DESKTOP, f"browse_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
            page.screenshot(path=shot_path, full_page=True)
            screenshot_msg = f"\nScreenshot saved: {shot_path}"

        # Extract text
        text = page.inner_text("body")[:6000]
        result_parts.append(screenshot_msg)
        result_parts.append(f"\nContent:\n{text}")
        return "\n".join(result_parts)

    # === Strategy 1: Connect to running Chrome via CDP ===
    try:
        with sync_playwright() as p:
            browser = p.chromium.connect_over_cdp(f"http://127.0.0.1:{CHROME_DEBUG_PORT}", timeout=3000)
            contexts = browser.contexts
            if contexts:
                page = contexts[0].new_page()
            else:
                page = browser.new_page()
            
            if stealth_available:
                stealth_sync(page)

            result = _process_page(page, url, action, wait, screenshot, "Chrome CDP â€” your real browser")
            page.close()
            # Don't close browser â€” it's the user's running Chrome
            return result
    except Exception:
        pass  # Chrome not running with debug port, try strategy 2

    # === Strategy 2: Launch with persistent profile ===
    # Load custom path from .env if set
    chrome_profile = CHROME_PROFILE_PATH
    env_path = os.path.join(AGENT_DIR, ".env")
    if os.path.exists(env_path):
        try:
            with open(env_path, 'r', encoding='utf-8') as f:
                for line in f:
                    line = line.strip()
                    if line.startswith('CHROME_PROFILE_PATH='):
                        chrome_profile = line.split('=', 1)[1].strip().strip('"').strip("'")
        except:
            pass

    if not os.path.exists(chrome_profile):
        return (f"Chrome profile not found at: {chrome_profile}\n\n"
                f"Fix: Add to .env file:\n"
                f"  CHROME_PROFILE_PATH=C:\\Users\\Dator\\AppData\\Local\\Google\\Chrome\\User Data\n\n"
                f"Or start Chrome with debug mode first:\n"
                f"  chrome.exe --remote-debugging-port={CHROME_DEBUG_PORT}")

    try:
        with sync_playwright() as p:
            # IMPORTANT: Chrome must be fully closed for this to work
            context = p.chromium.launch_persistent_context(
                user_data_dir=chrome_profile,
                headless=False,
                channel="chrome",  # Use installed Chrome, not Playwright's Chromium
                args=[
                    '--disable-blink-features=AutomationControlled',
                    '--no-first-run',
                    '--no-default-browser-check',
                    f'--remote-debugging-port={CHROME_DEBUG_PORT}'
                ],
                viewport={'width': 1920, 'height': 1080},
                locale='sv-SE',
                timezone_id='Europe/Stockholm',
                ignore_default_args=['--enable-automation']
            )

            page = context.new_page()
            if stealth_available:
                stealth_sync(page)

            result = _process_page(page, url, action, wait, screenshot, 
                                   "Chrome Profile â€” your passwords and logins active")
            context.close()
            return result

    except Exception as e:
        error_msg = str(e)
        if "already running" in error_msg.lower() or "lock" in error_msg.lower() or "user data" in error_msg.lower():
            return (f"Chrome is currently open. Two options:\n\n"
                    f"Option A â€” Close Chrome completely, then try again.\n\n"
                    f"Option B â€” Keep Chrome open but restart it with debug mode:\n"
                    f"  1. Close Chrome\n"
                    f"  2. Run in PowerShell:\n"
                    f"     & 'C:\\Program Files\\Google\\Chrome\\Application\\chrome.exe' --remote-debugging-port={CHROME_DEBUG_PORT}\n"
                    f"  3. Now browse_as_me will connect to your open Chrome\n\n"
                    f"Option B is best â€” Chrome stays open and agent can use it anytime.")
        return f"Error launching Chrome with profile: {error_msg}"


def tool_system_info() -> str:
    """Get system information"""
    code = """
import platform
try:
    import psutil
    print(f"=== System Info ===")
    print(f"OS: {platform.system()} {platform.release()} ({platform.version()})")
    print(f"Machine: {platform.machine()}")
    print(f"Processor: {platform.processor()}")
    print()
    print(f"=== CPU ===")
    print(f"Cores: {psutil.cpu_count(logical=False)} physical, {psutil.cpu_count()} logical")
    print(f"Usage: {psutil.cpu_percent(interval=1)}%")
    print()
    print(f"=== Memory ===")
    mem = psutil.virtual_memory()
    print(f"Total: {mem.total/(1024**3):.1f} GB")
    print(f"Used: {mem.used/(1024**3):.1f} GB ({mem.percent}%)")
    print(f"Free: {mem.available/(1024**3):.1f} GB")
    print()
    print(f"=== Disk ===")
    for part in psutil.disk_partitions():
        try:
            usage = psutil.disk_usage(part.mountpoint)
            print(f"{part.device}: {usage.used/(1024**3):.1f}/{usage.total/(1024**3):.1f} GB ({usage.percent}%)")
        except:
            pass
except ImportError:
    print("psutil not installed. Run: pip install psutil")
"""
    return tool_run_python(code)


def tool_clean_temp(dry_run: bool = True) -> str:
    """Clean temporary files"""
    temp_dirs = [
        os.environ.get('TEMP', ''),
        os.path.join(os.environ.get('LOCALAPPDATA', ''), 'Temp'),
    ]
    result = "=== Temp Files Cleanup ===\n"
    total_size = 0
    total_files = 0
    for temp_dir in temp_dirs:
        if not temp_dir or not os.path.exists(temp_dir):
            continue
        result += f"\n{temp_dir}:\n"
        dir_size = 0
        dir_files = 0
        for root, dirs, files in os.walk(temp_dir):
            for name in files:
                try:
                    fp = os.path.join(root, name)
                    size = os.path.getsize(fp)
                    dir_size += size
                    dir_files += 1
                    if not dry_run:
                        try:
                            os.remove(fp)
                        except (PermissionError, OSError):
                            pass
                except OSError:
                    continue
        result += f"  Files: {dir_files}, Size: {dir_size/(1024**2):.1f} MB\n"
        total_size += dir_size
        total_files += dir_files
    mode = "DRY RUN (no files deleted)" if dry_run else "CLEANED"
    result += f"\n--- {mode} ---\n"
    result += f"Total: {total_files} files, {total_size/(1024**2):.1f} MB\n"
    if dry_run:
        result += "\nTo actually delete, run with dry_run=false"
    return result


def tool_send_email(to: str, subject: str, body: str) -> str:
    """Send email via Gmail"""
    env_file = os.path.join(AGENT_DIR, ".env")
    if not os.path.exists(env_file):
        return """Email not configured yet. To set up Gmail:
1. Go to https://myaccount.google.com/apppasswords
2. Create an App Password for 'Mail'
3. Create file .env next to agent_v3.py with:
   GMAIL_ADDRESS=your.email@gmail.com
   GMAIL_APP_PASSWORD=your_app_password
4. Try again"""
    gmail_address = ""
    gmail_password = ""
    with open(env_file, 'r', encoding='utf-8', errors='replace') as f:
        for line in f:
            if line.startswith('GMAIL_ADDRESS='):
                gmail_address = line.split('=', 1)[1].strip()
            elif line.startswith('GMAIL_APP_PASSWORD='):
                gmail_password = line.split('=', 1)[1].strip()
    if not gmail_address or not gmail_password:
        return "Error: GMAIL_ADDRESS or GMAIL_APP_PASSWORD not found in .env file"
    try:
        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart
        msg = MIMEMultipart()
        msg['From'] = gmail_address
        msg['To'] = to
        msg['Subject'] = subject
        msg.attach(MIMEText(body, 'plain', 'utf-8'))
        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
            server.login(gmail_address, gmail_password)
            server.send_message(msg)
        return f"Email sent to {to} with subject '{subject}'"
    except Exception as e:
        return f"Error sending email: {e}"


# ============================================================
# FILE OPERATIONS (DELETE / MOVE)
# ============================================================

def _is_path_within_folder(path: str, folder: str) -> bool:
    """Return True if path is inside folder (or equal)."""
    try:
        norm_path = os.path.normpath(os.path.abspath(path))
        norm_folder = os.path.normpath(os.path.abspath(folder))
        if os.name == "nt":
            norm_path = os.path.normcase(norm_path)
            norm_folder = os.path.normcase(norm_folder)
        return os.path.commonpath([norm_path, norm_folder]) == norm_folder
    except Exception:
        return False


def _get_allowed_delete_roots(explicit_allowed_folder: str = None) -> list:
    roots = []
    if explicit_allowed_folder:
        roots.append(os.path.normpath(os.path.abspath(explicit_allowed_folder)))
    if LAST_SCANNED_FOLDER:
        roots.append(os.path.normpath(os.path.abspath(LAST_SCANNED_FOLDER)))
    return roots


def get_drive_root(path: str) -> str:
    """Return drive root for a path, e.g. E:\\ for Windows paths."""
    abs_path = os.path.abspath(path)
    drive, _ = os.path.splitdrive(abs_path)
    if drive:
        return drive + os.sep
    return os.path.sep


def get_relative_from_drive(path: str) -> str:
    """Return path relative to drive root."""
    abs_path = os.path.abspath(path)
    drive_root = get_drive_root(abs_path)
    rel = os.path.relpath(abs_path, drive_root)
    if rel in (".", ""):
        return ""
    return rel


def get_global_trash_path(path: str) -> str:
    """Build global per-drive trash destination, preserving source relative path."""
    drive_root = get_drive_root(path)
    relative_path = get_relative_from_drive(path)
    trash_root = os.path.join(drive_root, "_TRASH")
    if relative_path:
        return os.path.join(trash_root, relative_path)
    return trash_root


def _get_unique_trash_destination(destination_path: str) -> str:
    """Build a non-colliding destination path in global trash."""
    candidate = destination_path
    if not os.path.exists(candidate):
        return candidate
    parent = os.path.dirname(candidate)
    base_name = os.path.basename(candidate)
    stem, ext = os.path.splitext(base_name)
    suffix = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    return os.path.join(parent, f"{stem}_{suffix}_{uuid.uuid4().hex[:8]}{ext}")


def _move_to_trash(path: str, target_root: str = None) -> str:
    """Move file/folder to global per-drive trash and return new location."""
    dst = get_global_trash_path(path)
    os.makedirs(os.path.dirname(dst), exist_ok=True)
    dst = _get_unique_trash_destination(dst)
    shutil.move(path, dst)
    return dst


def _original_path_from_trash(trash_path: str) -> "str | None":
    """Reconstruct original path from a trash path.
    Trash structure: {drive}:\\_TRASH\\{relative}
    Original:        {drive}:\\{relative}
    Returns None if the path is not inside a _TRASH folder.
    """
    abs_path = os.path.abspath(trash_path)
    parts = abs_path.replace("\\", "/").split("/")
    try:
        trash_idx = next(i for i, p in enumerate(parts) if p == "_TRASH")
    except StopIteration:
        return None
    drive_parts = parts[:trash_idx]       # e.g. ['E:']
    rel_parts   = parts[trash_idx + 1:]   # relative path after _TRASH
    if not rel_parts:
        return None
    original = "/".join(drive_parts + rel_parts)
    # Normalise to OS separators
    return original.replace("/", os.sep)


def tool_restore_from_trash(trash_path: str) -> str:
    """Restore a file or folder from global _TRASH to its original location."""
    if not os.path.exists(trash_path):
        return f"Not found in trash: {trash_path}"
    original = _original_path_from_trash(trash_path)
    if not original:
        return f"Cannot determine original path for: {trash_path}"
    if os.path.exists(original):
        # Avoid overwriting: add suffix to original name
        stem, ext = os.path.splitext(original)
        suffix = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        original = f"{stem}_restored_{suffix}{ext}"
    os.makedirs(os.path.dirname(original), exist_ok=True)
    shutil.move(trash_path, original)
    return f"Restored: {trash_path} → {original}"


def tool_list_trash(drive: str = None) -> str:
    """List all items currently in the global _TRASH folder for the given drive (or all drives)."""
    # Detect drives to scan
    if drive:
        roots = [drive.rstrip("/\\") + os.sep]
    elif os.name == "nt":
        import string
        roots = [
            f"{d}:{os.sep}" for d in string.ascii_uppercase
            if os.path.isdir(f"{d}:{os.sep}")
        ]
    else:
        roots = [os.sep]

    entries = []
    total_size = 0
    for root in roots:
        trash_root = os.path.join(root, "_TRASH")
        if not os.path.isdir(trash_root):
            continue
        for dirpath, _dirs, files in os.walk(trash_root):
            for fname in files:
                full = os.path.join(dirpath, fname)
                try:
                    sz = os.path.getsize(full)
                    total_size += sz
                    entries.append((full, sz))
                except OSError:
                    pass

    if not entries:
        return "Trash is empty."
    lines = [f"  {path}  ({sz/1024:.1f} KB)" for path, sz in entries[:50]]
    result = f"Trash contents ({len(entries)} file(s), {total_size/1024/1024:.1f} MB total):\n"
    result += "\n".join(lines)
    if len(entries) > 50:
        result += f"\n  … and {len(entries)-50} more"
    return result


def tool_purge_trash(drive: str = None, confirm: bool = False) -> str:
    """Permanently delete everything in the global _TRASH folder.
    Requires confirm=True to actually delete (safety guard).
    """
    if not confirm:
        return (
            "Purge requires confirm=True. "
            "This will permanently delete all items in _TRASH — irreversible!"
        )
    if drive:
        roots = [drive.rstrip("/\\") + os.sep]
    elif os.name == "nt":
        import string
        roots = [
            f"{d}:{os.sep}" for d in string.ascii_uppercase
            if os.path.isdir(f"{d}:{os.sep}")
        ]
    else:
        roots = [os.sep]

    purged = 0
    errors = []
    for root in roots:
        trash_root = os.path.join(root, "_TRASH")
        if not os.path.isdir(trash_root):
            continue
        for item in os.listdir(trash_root):
            full = os.path.join(trash_root, item)
            try:
                if os.path.isfile(full) or os.path.islink(full):
                    os.remove(full)
                    purged += 1
                elif os.path.isdir(full):
                    shutil.rmtree(full)
                    purged += 1
            except Exception as e:
                errors.append(f"{full}: {e}")

    result = f"Purged {purged} item(s) from _TRASH permanently."
    if errors:
        result += f"\nErrors ({len(errors)}):\n" + "\n".join(errors)
    return result


def tool_clean_duplicates(path: str, mode: str = "trash", keep: str = "newest", permanent: bool = False, confirm: bool = False, allowed_folder: str = None) -> str:
    """Find duplicates and move removable copies to global trash (reversible by default)."""
    if not path:
        return "Path is required."
    if not os.path.isdir(path):
        return f"Folder not found: {path}"

    global LAST_SCANNED_FOLDER
    LAST_SCANNED_FOLDER = os.path.abspath(path)

    files_by_key = {}
    for root, dirs, files in os.walk(path):
        for name in files:
            fp = os.path.join(root, name)
            try:
                size = os.path.getsize(fp)
                key = (name, size)
                files_by_key.setdefault(key, []).append(fp)
            except OSError:
                continue

    duplicates = {k: v for k, v in files_by_key.items() if len(v) > 1}
    if not duplicates:
        return f"No duplicates found in {path}"

    to_remove = []
    kept = 0
    for _, group in duplicates.items():
        group_sorted = sorted(
            group,
            key=lambda p: os.path.getmtime(p),
            reverse=(keep == "newest")
        )
        kept += 1
        to_remove.extend(group_sorted[1:])

    if mode == "dry_run":
        total_bytes = sum(os.path.getsize(p) for p in to_remove if os.path.exists(p))
        mb = total_bytes / (1024 * 1024)
        size_str = f"{mb:.1f} MB" if mb >= 1 else f"{total_bytes/1024:.1f} KB"
        return (
            f"Dry run: found {len(duplicates)} duplicate groups.\n"
            f"Would move {len(to_remove)} file(s) ({size_str}) to global trash (reversible)."
        )

    if mode == "permanent":
        permanent = True

    delete_result = tool_delete_files(
        to_remove,
        permanent=permanent,
        confirm=confirm,
        allowed_folder=allowed_folder or path,
    )
    if delete_result.startswith("Blocked"):
        return delete_result

    action = "permanently deleted" if permanent else "moved to global trash (reversible)"
    return (
        f"Cleaned {len(to_remove)} duplicate file(s); kept {kept} original(s).\n"
        f"Files {action}.\n"
        f"{delete_result}"
    )


def tool_delete_files(paths: list, permanent: bool = False, confirm: bool = False, allowed_folder: str = None) -> str:
    """Delete one or more files. Safe mode moves to global trash by default."""
    if not paths:
        return "No file paths provided."

    if permanent and not confirm:
        return "Blocked: permanent delete requires permanent=true and confirm=true."

    moved = []
    errors = []
    allowed_roots = _get_allowed_delete_roots(allowed_folder)

    for path in paths:
        try:
            path = path.strip()
            if not path:
                continue

            matched_root = None
            for root in allowed_roots:
                if _is_path_within_folder(path, root):
                    matched_root = root
                    break
            if not matched_root:
                errors.append(f"  Blocked (outside allowed folders): {path}")
                continue

            if os.path.isfile(path):
                size = os.path.getsize(path)
                if permanent:
                    os.remove(path)
                    new_location = "(permanent)"
                else:
                    new_location = _move_to_trash(path, matched_root)
                if size < 1024:
                    size_str = f"{size} B"
                elif size < 1024*1024:
                    size_str = f"{size/1024:.1f} KB"
                else:
                    size_str = f"{size/(1024*1024):.1f} MB"
                if permanent:
                    moved.append(f"  Permanently deleted: {os.path.basename(path)} ({size_str})")
                else:
                    moved.append(
                        f"  Moved to global trash:\n"
                        f"    {path}\n"
                        f"    {new_location}\n"
                        f"    ({size_str})"
                    )
            elif os.path.isdir(path):
                if permanent:
                    shutil.rmtree(path)
                    moved.append(f"  Permanently deleted folder: {path}")
                else:
                    new_location = _move_to_trash(path, matched_root)
                    moved.append(f"  Moved to global trash:\n    {path}\n    {new_location}")
            else:
                errors.append(f"  Not found: {path}")
        except PermissionError:
            errors.append(f"  Permission denied: {path}")
        except Exception as e:
            errors.append(f"  Error deleting {path}: {e}")

    result = ""
    if moved:
        verb = "deleted" if permanent else "moved"
        result += f"Successfully {verb} {len(moved)} item(s):\n" + "\n".join(moved)
    if errors:
        result += f"\n\nErrors ({len(errors)}):\n" + "\n".join(errors)

    return result if result else "Nothing to delete."


def tool_move_files(files: list) -> str:
    """Move files from one location to another"""
    if not files:
        return "No files provided."
    
    moved = []
    errors = []
    
    for item in files:
        src = item.get("from", "")
        dst = item.get("to", "")
        try:
            if not os.path.exists(src):
                errors.append(f"  Not found: {src}")
                continue
            os.makedirs(os.path.dirname(dst), exist_ok=True)
            shutil.move(src, dst)
            moved.append(f"  {os.path.basename(src)} -> {dst}")
        except Exception as e:
            errors.append(f"  Error moving {src}: {e}")
    
    result = ""
    if moved:
        result += f"Moved {len(moved)} file(s):\n" + "\n".join(moved)
    if errors:
        result += f"\n\nErrors ({len(errors)}):\n" + "\n".join(errors)
    
    return result if result else "Nothing to move."


# ============================================================
# MEMORY TOOLS
# ============================================================

def tool_memory_add_fact(fact: str) -> str:
    """Add a fact about the user to persistent memory"""
    profile = load_user_profile()
    facts = profile.get("facts", [])
    if fact not in facts:
        facts.append(fact)
        profile["facts"] = facts[-50:]  # Keep last 50 facts
        save_user_profile(profile)
        return f"Remembered: {fact}"
    return f"Already known: {fact}"


def tool_memory_add_task(task: str) -> str:
    """Add a task to the persistent task list"""
    tasks = load_tasks()
    new_task = {
        "task": task,
        "status": "pending",
        "added": datetime.datetime.now().isoformat(),
    }
    tasks.append(new_task)
    save_tasks(tasks)
    return f"Task added: {task} (#{len(tasks)-1})"


def tool_memory_complete_task(task_index: int) -> str:
    """Mark a task as complete"""
    tasks = load_tasks()
    if 0 <= task_index < len(tasks):
        tasks[task_index]["status"] = "done"
        tasks[task_index]["completed"] = datetime.datetime.now().isoformat()
        save_tasks(tasks)
        return f"Task #{task_index} marked as done: {tasks[task_index]['task']}"
    return f"Task #{task_index} not found. Total tasks: {len(tasks)}"


def tool_memory_show() -> str:
    """Show current memory state"""
    return get_memory_context()


# ============================================================
# N8N INTEGRATION
# ============================================================

def _get_n8n_api_key() -> str:
    """Read n8n API key from environment or .env file."""
    env_key = os.environ.get("N8N_API_KEY", "").strip()
    if env_key:
        return env_key

    env_file = os.path.join(AGENT_DIR, ".env")
    if os.path.exists(env_file):
        with open(env_file, 'r', encoding='utf-8', errors='replace') as f:
            for line in f:
                if line.startswith('N8N_API_KEY='):
                    return line.split('=', 1)[1].strip()
    return ""


def _get_n8n_base_url() -> str:
    """Read n8n base URL from env or .env, fallback to localhost."""
    env_url = os.environ.get("N8N_URL", "").strip()
    if env_url:
        return env_url.rstrip("/")

    env_file = os.path.join(AGENT_DIR, ".env")
    if os.path.exists(env_file):
        with open(env_file, 'r', encoding='utf-8', errors='replace') as f:
            for line in f:
                if line.startswith('N8N_URL='):
                    url = line.split('=', 1)[1].strip()
                    if url:
                        return url.rstrip("/")

    return N8N_URL.rstrip("/")


def _n8n_request(method: str, path: str, data: dict = None, params: dict = None, timeout: int = 30) -> dict:
    """Make authenticated request to n8n API."""
    api_key = _get_n8n_api_key()
    if not api_key:
        return {"error": "n8n API key not configured. Add N8N_API_KEY=your_key to .env file next to agent_v3.py"}

    headers = {
        "accept": "application/json",
        "content-type": "application/json",
        "X-N8N-API-KEY": api_key
    }

    base_url = _get_n8n_base_url()
    url = f"{base_url}/api/v1{path}"
    try:
        if method == "GET":
            resp = requests.get(url, headers=headers, params=params, timeout=timeout)
        elif method == "POST":
            resp = requests.post(url, headers=headers, json=data, params=params, timeout=timeout)
        elif method == "PUT":
            resp = requests.put(url, headers=headers, json=data, params=params, timeout=timeout)
        elif method == "PATCH":
            resp = requests.patch(url, headers=headers, json=data, params=params, timeout=timeout)
        elif method == "DELETE":
            resp = requests.delete(url, headers=headers, params=params, timeout=timeout)
        else:
            return {"error": f"Unknown method: {method}"}

        if resp.status_code >= 400:
            return {
                "error": f"n8n API error {resp.status_code}: {resp.text[:500]}",
                "status_code": resp.status_code,
                "url": url,
            }

        if not resp.text:
            return {}
        try:
            return resp.json()
        except Exception:
            return {"raw": resp.text}
    except requests.exceptions.ConnectionError:
        return {"error": f"Cannot connect to n8n at {base_url}. Is it running? Start with: n8n start"}
    except Exception as e:
        return {"error": f"n8n request failed: {e}"}


def _n8n_find_workflow_by_name(workflow_name: str) -> dict:
    """Resolve workflow by exact/fuzzy name."""
    if not workflow_name:
        return {"error": "workflow_name is required"}

    listing = _n8n_request("GET", "/workflows", params={"limit": 200}, timeout=20)
    if "error" in listing:
        return listing

    workflows = listing.get("data", [])
    exact = [w for w in workflows if str(w.get("name", "")).strip().lower() == workflow_name.strip().lower()]
    if exact:
        return _n8n_request("GET", f"/workflows/{exact[0].get('id')}")

    fuzzy = [w for w in workflows if workflow_name.strip().lower() in str(w.get("name", "")).strip().lower()]
    if len(fuzzy) == 1:
        return _n8n_request("GET", f"/workflows/{fuzzy[0].get('id')}")

    if len(fuzzy) > 1:
        return {
            "error": f"Multiple workflows matched '{workflow_name}'",
            "matches": [{"id": w.get("id"), "name": w.get("name"), "active": bool(w.get("active"))} for w in fuzzy[:10]]
        }

    return {"error": f"Workflow '{workflow_name}' not found"}


def tool_n8n_list_workflows(query: str = "", raw: bool = False) -> str:
    """List n8n workflows (id + name + active)."""
    result = _n8n_request("GET", "/workflows", params={"limit": 200}, timeout=20)
    if "error" in result:
        return result["error"]

    workflows = []
    for wf in result.get("data", []):
        row = {
            "id": wf.get("id"),
            "name": wf.get("name", "Unnamed"),
            "active": bool(wf.get("active")),
            "updatedAt": wf.get("updatedAt"),
        }
        if query and query.lower() not in str(row["name"]).lower():
            continue
        workflows.append(row)

    if raw:
        return json.dumps({"data": workflows}, ensure_ascii=False)

    if not workflows:
        return "No workflows found in n8n."

    output = f"Found {len(workflows)} workflow(s):\n\n"
    for wf in workflows:
        status = "ACTIVE" if wf.get("active") else "inactive"
        name = wf.get("name", "Unnamed")
        wf_id = wf.get("id", "?")
        updated = str(wf.get("updatedAt", ""))[:10]
        output += f"  [{status}] {name}\n"
        output += f"    ID: {wf_id} | Updated: {updated}\n\n"

    return output


def tool_n8n_get_workflow(wf_id: str = "", name: str = "", raw: bool = False) -> str:
    """Get full workflow JSON by id or name."""
    if wf_id:
        result = _n8n_request("GET", f"/workflows/{wf_id}")
    else:
        result = _n8n_find_workflow_by_name(name)
    if "error" in result:
        return json.dumps(result, ensure_ascii=False) if raw else result["error"]

    if raw:
        return json.dumps(result, ensure_ascii=False)

    wf = result
    wf_name = wf.get("name", "Unnamed")
    active = "ACTIVE" if wf.get("active") else "inactive"
    nodes = wf.get("nodes", [])
    connections = wf.get("connections", {})

    output = f"Workflow: {wf_name}\n"
    output += f"ID: {wf.get('id')} | Status: {active}\n"
    output += f"Created: {wf.get('createdAt', '')[:10]} | Updated: {wf.get('updatedAt', '')[:10]}\n\n"
    output += f"Nodes ({len(nodes)}):\n"
    for node in nodes:
        output += f"  - {node.get('name', '?')} ({node.get('type', '?')})\n"

    output += f"\nConnections: {len(connections)} link(s)\n"
    return output


def tool_n8n_get_executions(workflow_id: str, limit: int = 5, raw: bool = False) -> str:
    """List recent executions for a workflow."""
    params = {"workflowId": workflow_id, "limit": max(1, min(int(limit or 5), 50))}
    result = _n8n_request("GET", "/executions", params=params, timeout=20)
    if "error" in result:
        return json.dumps(result, ensure_ascii=False) if raw else result["error"]

    executions = []
    for ex in result.get("data", []):
        executions.append({
            "id": ex.get("id"),
            "status": ex.get("status") or ("success" if ex.get("finished") and not ex.get("stoppedAt") else "error"),
            "startedAt": ex.get("startedAt"),
        })

    if raw:
        return json.dumps({"data": executions}, ensure_ascii=False)

    if not executions:
        return f"No executions found for workflow {workflow_id}."
    lines = [f"Executions for workflow {workflow_id}:"]
    for ex in executions:
        lines.append(f"- id={ex['id']} status={ex['status']} startedAt={ex['startedAt']}")
    return "\n".join(lines)


def tool_n8n_get_execution(execution_id: str, raw: bool = False) -> str:
    """Get full execution details including errors/runData."""
    result = _n8n_request("GET", f"/executions/{execution_id}", params={"includeData": "true"}, timeout=30)
    if "error" in result:
        return json.dumps(result, ensure_ascii=False) if raw else result["error"]
    if raw:
        return json.dumps(result, ensure_ascii=False)

    status = result.get("status", "unknown")
    started = result.get("startedAt", "")
    last_node = result.get("data", {}).get("resultData", {}).get("lastNodeExecuted")
    err = result.get("data", {}).get("resultData", {}).get("error", {})
    err_msg = err.get("message", "") if isinstance(err, dict) else str(err)
    return f"Execution {execution_id}: status={status}, startedAt={started}, lastNode={last_node}, error={err_msg[:500]}"


def _extract_execution_id(run_result: dict) -> str:
    """Handle multiple n8n response shapes for run endpoint."""
    if not isinstance(run_result, dict):
        return ""
    for key in ("executionId", "id"):
        if run_result.get(key):
            return str(run_result[key])
    data = run_result.get("data")
    if isinstance(data, dict):
        for key in ("executionId", "id"):
            if data.get(key):
                return str(data[key])
    return ""


def _is_n8n_http_error(result: dict, codes: tuple) -> bool:
    if not isinstance(result, dict) or "error" not in result:
        return False
    text = str(result.get("error", "")).lower()
    return any(f" {code}" in text or f"{code}:" in text for code in codes)


# ---------------------------------------------------------------------------
# Webhook-based workflow runner (bypasses 405 on POST /workflows/{id}/run)
# ---------------------------------------------------------------------------

_AGENT_WEBHOOK_PATH = "agent-auto-run"  # custom path used for all auto-runs


def _n8n_get_webhook_path(workflow_id: str) -> str:
    """Return the DB webhookPath for the given workflow, or '' if not found.

    n8n stores webhook registrations in webhook_entity with path format:
      {workflowId}/webhook/{customPath}
    The production URL is: {N8N_URL}/webhook/{webhookPath}
    """
    try:
        import sqlite3 as _sqlite3
        db_path = os.path.join(os.path.expanduser("~"), ".n8n", "database.sqlite")
        conn = _sqlite3.connect(db_path, timeout=5)
        cur = conn.cursor()
        cur.execute(
            "SELECT webhookPath FROM webhook_entity WHERE workflowId=? LIMIT 1",
            (workflow_id,),
        )
        row = cur.fetchone()
        conn.close()
        return row[0] if row else ""
    except Exception:
        return ""


def _n8n_add_webhook_trigger(workflow_id: str, workflow_json: dict) -> str:
    """Add / replace a webhook trigger so the workflow can be fired on demand.

    Strategy:
      - manualTrigger present  → replace it with the Webhook node (cleaner).
      - scheduleTrigger present → keep it; add Webhook alongside with copied connections
        so both entry-points fire the same chain.
      - No recognised trigger  → prepend Webhook unconnected (rare fallback).

    Returns the webhookPath string ('{workflowId}/webhook/{customPath}') or '' on failure.
    The last error detail is stored in _n8n_add_webhook_trigger.last_error (str).
    """
    import copy as _copy, uuid as _uuid

    _n8n_add_webhook_trigger.last_error = ""

    nodes = list(workflow_json.get("nodes", []))
    connections = dict(workflow_json.get("connections", {}))

    webhook_node_id = f"agent-wh-{_uuid.uuid4().hex[:8]}"
    custom_path = _AGENT_WEBHOOK_PATH

    # Use "Webhook" (no spaces) so the DB path stays "{workflowId}/webhook/{customPath}"
    wh_node_name = "Webhook"
    webhook_node = {
        "parameters": {
            "httpMethod": "POST",
            "path": custom_path,
            "responseMode": "onReceived",
            "options": {},
        },
        "id": webhook_node_id,
        "name": wh_node_name,
        "type": "n8n-nodes-base.webhook",
        "typeVersion": 2,
        "position": [0, 300],
    }

    # Locate existing trigger nodes (by type substring)
    _REPLACEABLE_TRIGGERS = ("manualTrigger", "scheduleTrigger")
    trigger_idx = next(
        (
            i
            for i, n in enumerate(nodes)
            if any(t in n.get("type", "") for t in _REPLACEABLE_TRIGGERS)
        ),
        None,
    )

    if trigger_idx is not None:
        # Replace the existing trigger with Webhook.
        # manualTrigger  → clean replacement, no side effects.
        # scheduleTrigger → n8n v2.7.4 cannot activate schedule-trigger workflows due to a
        #   locale bug ("Unknown alias: und"); replacing with Webhook lets us activate and
        #   test the workflow logic.  The schedule can be restored via n8n UI afterwards.
        old_name = nodes[trigger_idx]["name"]
        nodes[trigger_idx] = webhook_node
        if old_name in connections:
            connections[wh_node_name] = connections.pop(old_name)
    else:
        # No recognised trigger — prepend Webhook (chain not auto-wired)
        nodes.insert(0, webhook_node)

    payload = _sanitize_workflow_payload(dict(workflow_json))
    payload["nodes"] = nodes
    payload["connections"] = connections

    result = _n8n_request("PUT", f"/workflows/{workflow_id}", data=payload, timeout=45)
    if "error" in result:
        err = result.get("error", str(result))
        _n8n_add_webhook_trigger.last_error = f"PUT /workflows/{workflow_id} failed: {err}"
        return ""

    # Activate so webhook registers in DB
    act = _n8n_request("POST", f"/workflows/{workflow_id}/activate", timeout=15)
    if "error" in act:
        err = act.get("error", str(act))
        _n8n_add_webhook_trigger.last_error = f"Activate failed: {err}"
        return ""

    # Give n8n a moment to register the webhook
    import time as _time
    _time.sleep(1.5)

    # Read the actual webhookPath from DB (n8n encodes node name in the path)
    db_path = _n8n_get_webhook_path(workflow_id)
    if not db_path:
        _n8n_add_webhook_trigger.last_error = (
            "Webhook registered in n8n but path not found in webhook_entity DB"
        )
    return db_path


_n8n_add_webhook_trigger.last_error = ""  # type: ignore[attr-defined]


def _n8n_trigger_via_webhook(webhook_path: str, timeout: int = 60) -> dict:
    """POST to production webhook URL and return parsed JSON body."""
    base = _get_n8n_base_url()
    url = f"{base}/webhook/{webhook_path}"
    try:
        resp = requests.post(url, json={}, timeout=timeout)
        if resp.status_code >= 400:
            return {"error": f"webhook {resp.status_code}: {resp.text[:300]}"}
        try:
            return resp.json()
        except Exception:
            return {"raw": resp.text}
    except Exception as e:
        return {"error": str(e)}


# Fields allowed in POST /workflows and PUT /workflows/{id} payloads.
# n8n rejects any field not in this whitelist with 400 "must NOT have additional properties".
_WORKFLOW_ALLOWED_FIELDS = frozenset({
    "name", "nodes", "connections", "settings", "staticData", "pinData",
})


_N8N_MAX_NAME_LEN = 128  # n8n enforces 1–128 chars on PUT /workflows/{id}


def _sanitize_workflow_payload(workflow_json: dict, force_name: str = "") -> dict:
    """Keep only fields that n8n accepts in POST/PUT workflow payloads.

    n8n rejects:
    - Any key not in _WORKFLOW_ALLOWED_FIELDS → 400 "must NOT have additional properties"
    - null values for optional fields (pinData, staticData) → 400
    - name longer than 128 chars → 400 "Workflow name must be 1 to 128 characters long"
    """
    payload = {
        k: v
        for k, v in (workflow_json or {}).items()
        if k in _WORKFLOW_ALLOWED_FIELDS and v is not None
    }
    if force_name and not payload.get("name"):
        payload["name"] = force_name
    # Truncate name to n8n's 128-char limit (preserves suffix "..." for readability)
    name = payload.get("name", "")
    if len(name) > _N8N_MAX_NAME_LEN:
        payload["name"] = name[: _N8N_MAX_NAME_LEN - 3] + "..."
    return payload


def tool_n8n_run_workflow(workflow_id: str, wait: bool = True, raw: bool = False) -> str:
    """Trigger workflow execution via webhook (bypasses 405 on POST /run).

    Strategy:
      1. Check if workflow already has a webhook registered (webhook_entity table).
      2. If not: patch the workflow to add/replace a manualTrigger with an Agent
         Webhook node, then activate.
      3. POST to the production webhook URL.
      4. Get the latest execution_id from GET /executions.
    """
    import time as _time

    # Step 1: Ensure webhook exists
    wh_path = _n8n_get_webhook_path(workflow_id)

    if not wh_path:
        # Need to add webhook trigger to workflow
        wf = _n8n_request("GET", f"/workflows/{workflow_id}", timeout=20)
        if "error" in wf:
            out = {"error": wf["error"]}
            return json.dumps(out, ensure_ascii=False) if raw else out["error"]
        wh_path = _n8n_add_webhook_trigger(workflow_id, wf)
        if not wh_path:
            detail = _n8n_add_webhook_trigger.last_error or "unknown reason"
            out = {"error": f"Failed to add webhook trigger to workflow: {detail}"}
            return json.dumps(out, ensure_ascii=False) if raw else out["error"]
    else:
        # Ensure workflow is active (webhook only fires if active)
        _n8n_request("POST", f"/workflows/{workflow_id}/activate", timeout=15)
        _time.sleep(1)

    # Step 2: Trigger via webhook
    wh_result = _n8n_trigger_via_webhook(wh_path, timeout=60 if wait else 10)
    if "error" in wh_result:
        out = {"error": wh_result["error"]}
        return json.dumps(out, ensure_ascii=False) if raw else out["error"]

    # Step 3: Get the latest execution_id
    _time.sleep(0.5)  # brief wait for execution to be recorded
    ex_list = _n8n_request(
        "GET", "/executions",
        params={"workflowId": workflow_id, "limit": 1},
        timeout=20,
    )
    execution_id = ""
    if isinstance(ex_list, dict):
        data = ex_list.get("data", [])
        if data:
            execution_id = str(data[0].get("id", ""))

    response = {"execution_id": execution_id, "webhook_triggered": True, "raw": wh_result}
    return json.dumps(response, ensure_ascii=False) if raw else f"Workflow run triggered via webhook. execution_id={execution_id or '?'}"


def tool_n8n_update_workflow(workflow_id: str, workflow_json: dict, raw: bool = False) -> str:
    """Update workflow JSON in n8n."""
    if not isinstance(workflow_json, dict):
        err = {"error": "workflow_json must be an object"}
        return json.dumps(err, ensure_ascii=False) if raw else err["error"]

    payload = _sanitize_workflow_payload(workflow_json)
    result = _n8n_request("PUT", f"/workflows/{workflow_id}", data=payload, timeout=45)
    if "error" in result:
        return json.dumps(result, ensure_ascii=False) if raw else result["error"]
    return json.dumps(result, ensure_ascii=False) if raw else f"Workflow {workflow_id} updated."


def tool_n8n_validate_workflow(workflow_json: dict, raw: bool = False) -> str:
    """Validate workflow structure locally."""
    errors = []
    if not isinstance(workflow_json, dict):
        errors.append("workflow_json must be an object")
    else:
        if not isinstance(workflow_json.get("nodes"), list) or not workflow_json.get("nodes"):
            errors.append("nodes must be a non-empty list")
        if not isinstance(workflow_json.get("connections"), dict):
            errors.append("connections must be an object")

        node_names = set()
        node_ids = set()
        for idx, node in enumerate(workflow_json.get("nodes", [])):
            name = node.get("name")
            node_id = node.get("id")
            if not name:
                errors.append(f"node[{idx}] missing name")
            else:
                node_names.add(name)
            if not node_id:
                errors.append(f"node[{idx}] missing id")
            elif node_id in node_ids:
                errors.append(f"duplicate node id: {node_id}")
            else:
                node_ids.add(node_id)

        for source_name, src_links in workflow_json.get("connections", {}).items():
            if source_name not in node_names:
                errors.append(f"connection source '{source_name}' has no matching node")
            if not isinstance(src_links, dict):
                errors.append(f"connection block '{source_name}' must be object")
                continue
            for edge_groups in src_links.values():
                if not isinstance(edge_groups, list):
                    continue
                for edge_group in edge_groups:
                    if not isinstance(edge_group, list):
                        continue
                    for edge in edge_group:
                        if not isinstance(edge, dict):
                            continue
                        target = edge.get("node")
                        if target and target not in node_names:
                            errors.append(f"connection target '{target}' has no matching node")

    result = {"valid": not errors, "errors": errors}
    return json.dumps(result, ensure_ascii=False) if raw else ("Workflow is valid." if not errors else "Validation failed: " + "; ".join(errors[:5]))


def tool_n8n_create_workflow(name: str = "", description: str = "", workflow_json: dict = None, raw: bool = False) -> str:
    """Create a new n8n workflow via provided JSON or Ollama generation fallback."""

    # Deterministic path: caller provides full workflow JSON.
    if isinstance(workflow_json, dict):
        payload = _sanitize_workflow_payload(workflow_json, force_name=name)
        if not payload.get("name"):
            err = {"error": "workflow_json.name is required"}
            return json.dumps(err, ensure_ascii=False) if raw else err["error"]
        if "nodes" not in payload or "connections" not in payload:
            err = {"error": "workflow_json must include nodes and connections"}
            return json.dumps(err, ensure_ascii=False) if raw else err["error"]

        result = _n8n_request("POST", "/workflows", payload)
        if "error" in result:
            return json.dumps(result, ensure_ascii=False) if raw else result["error"]

        if raw:
            return json.dumps(result, ensure_ascii=False)
        wf_id = result.get("id", "?")
        return f"Workflow created successfully: {payload.get('name')} (ID: {wf_id})"
    
    # First, ask Ollama to generate the n8n workflow JSON
    generation_prompt = f"""Generate a complete n8n workflow JSON for this task:
Name: {name}
Description: {description}

Return ONLY valid JSON (no markdown, no explanation) with this structure:
{{
    "name": "{name}",
    "nodes": [
        {{
            "parameters": {{}},
            "type": "n8n-nodes-base.scheduleTrigger",
            "typeVersion": 1.2,
            "position": [250, 300],
            "id": "unique-uuid-here",
            "name": "Schedule Trigger"
        }}
    ],
    "connections": {{}},
    "settings": {{
        "executionOrder": "v1"
    }}
}}

Common node types:
- n8n-nodes-base.scheduleTrigger (cron/interval triggers)
- n8n-nodes-base.httpRequest (HTTP calls)
- n8n-nodes-base.code (JavaScript code)
- n8n-nodes-base.gmail / n8n-nodes-base.gmailTrigger
- n8n-nodes-base.googleSheets (read/write spreadsheets)
- n8n-nodes-base.telegram (send messages)
- n8n-nodes-base.if (conditions)
- n8n-nodes-base.set (set variables)
- n8n-nodes-base.merge (merge data)
- n8n-nodes-base.splitInBatches (loop over items)
- @n8n/n8n-nodes-langchain.lmChatOllama (local Ollama LLM)
- @n8n/n8n-nodes-langchain.agent (AI agent)

For scheduleTrigger with cron, use:
"parameters": {{ "rule": {{ "interval": [{{ "triggerAtHour": 9 }}] }} }}

For httpRequest:
"parameters": {{ "url": "https://example.com", "method": "GET" }}

For code node:
"parameters": {{ "jsCode": "// your code here\\nreturn items;" }}

IMPORTANT: 
- Each node MUST have a unique "id" (use format like "node-1", "node-2", etc.)
- Connections format: {{ "Node Name": {{ "main": [[{{ "node": "Next Node Name", "type": "main", "index": 0 }}]] }} }}
- Position each node 200px apart horizontally
- Return ONLY the JSON, nothing else."""

    try:
        resp = requests.post(f"{OLLAMA_URL}/api/chat", json={
            "model": MODEL,
            "messages": [
                {"role": "system", "content": "You are an n8n workflow generator. Return ONLY valid JSON. No markdown, no explanation, no code blocks."},
                {"role": "user", "content": generation_prompt}
            ],
            "stream": False,
            "options": {"temperature": 0.2, "num_predict": 8192}
        }, timeout=180)

        answer = resp.json()["message"]["content"].strip()

        # Clean up response - extract JSON
        if "```json" in answer:
            answer = answer.split("```json")[1].split("```")[0].strip()
        elif "```" in answer:
            answer = answer.split("```")[1].split("```")[0].strip()

        # Find JSON boundaries
        start = answer.find("{")
        end = answer.rfind("}") + 1
        if start == -1 or end <= start:
            return f"Error: Ollama did not generate valid JSON. Raw response:\n{answer[:500]}"

        workflow_json = json.loads(answer[start:end])

    except json.JSONDecodeError as e:
        return f"Error: Could not parse generated workflow JSON: {e}\nRaw:\n{answer[:500]}"
    except Exception as e:
        return f"Error generating workflow with Ollama: {e}"

    # Ensure name is set
    workflow_json["name"] = name
    workflow_json = _sanitize_workflow_payload(workflow_json)

    # Create the workflow in n8n
    result = _n8n_request("POST", "/workflows", workflow_json)
    if "error" in result:
        return f"Workflow generated but failed to create in n8n: {result['error']}\n\nGenerated JSON:\n{json.dumps(workflow_json, indent=2)[:2000]}"

    wf_id = result.get("id", "?")
    nodes_count = len(result.get("nodes", []))

    output = f"Workflow created successfully!\n"
    output += f"  Name: {name}\n"
    output += f"  ID: {wf_id}\n"
    output += f"  Nodes: {nodes_count}\n"
    output += f"  Status: inactive (activate with n8n_activate_workflow)\n"
    output += f"\n  Open in n8n: http://localhost:5678/workflow/{wf_id}\n"
    output += f"\n  To activate: use n8n_activate_workflow with id={wf_id}"

    return output


def tool_n8n_activate_workflow(wf_id: str, active: bool = True) -> str:
    """Activate or deactivate a workflow"""
    result = _n8n_request("PATCH", f"/workflows/{wf_id}", {"active": active})
    if "error" in result:
        return result["error"]

    status = "activated" if active else "deactivated"
    name = result.get("name", "Unnamed")
    return f"Workflow '{name}' (ID: {wf_id}) {status} successfully."


def tool_n8n_delete_workflow(wf_id: str) -> str:
    """Delete a workflow"""
    info = _n8n_request("GET", f"/workflows/{wf_id}")
    name = info.get("name", "Unknown") if "error" not in info else "Unknown"

    result = _n8n_request("DELETE", f"/workflows/{wf_id}")
    if "error" in result:
        return result["error"]

    return f"Workflow '{name}' (ID: {wf_id}) deleted."

# ============================================================
# SKILLS LOADER
# ============================================================

LOADED_SKILLS = {}  # skill_name -> module
SKILL_TOOLS_MAP = {}  # tool_name -> callable

def load_skills():
    """Load all enabled skills from the skills/ directory"""
    global LOADED_SKILLS, SKILL_TOOLS_MAP
    
    if not os.path.exists(SKILLS_DIR):
        os.makedirs(SKILLS_DIR, exist_ok=True)
        print("  Skills: created skills/ directory")
        return
    
    # Load whitelist
    whitelist_path = os.path.join(SKILLS_DIR, "skills_enabled.json")
    enabled = None
    if os.path.exists(whitelist_path):
        try:
            with open(whitelist_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                enabled = data.get("enabled", [])
        except:
            pass
    
    # Scan for skill files
    import importlib.util
    skill_files = [f for f in os.listdir(SKILLS_DIR) 
                   if f.endswith('.py') and not f.startswith('_')]
    
    loaded_count = 0
    total_tools = 0
    
    for skill_file in sorted(skill_files):
        skill_name = skill_file[:-3]  # Remove .py
        
        # Check whitelist
        if enabled is not None and skill_name not in enabled:
            continue
        
        try:
            spec = importlib.util.spec_from_file_location(
                f"skills.{skill_name}",
                os.path.join(SKILLS_DIR, skill_file)
            )
            module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(module)
            
            # Check if it has TOOLS dict
            if hasattr(module, 'TOOLS') and isinstance(module.TOOLS, dict):
                LOADED_SKILLS[skill_name] = module
                for tool_name, tool_func in module.TOOLS.items():
                    SKILL_TOOLS_MAP[tool_name] = tool_func
                    total_tools += 1
                loaded_count += 1
        except Exception as e:
            print(f"  Skills: error loading {skill_name}: {e}")
    
    if loaded_count > 0:
        print(f"  Skills: loaded {loaded_count} skills with {total_tools} tools")
    else:
        print(f"  Skills: no skills found in {SKILLS_DIR}")


def get_skills_description() -> str:
    """Generate tool descriptions for all loaded skills"""
    if not LOADED_SKILLS:
        return ""
    
    parts = ["\n\n=== SKILLS (loaded from skills/ directory) ===\n"]
    tool_num = 30  # Start numbering after built-in tools
    
    for skill_name, module in sorted(LOADED_SKILLS.items()):
        skill_desc = getattr(module, 'SKILL_DESCRIPTION', skill_name)
        skill_tools = getattr(module, 'SKILL_TOOLS', {})
        
        parts.append(f"\n[{skill_name}] â€” {skill_desc}")
        
        for tool_name, tool_info in skill_tools.items():
            desc = tool_info.get('description', '')
            example = tool_info.get('example', '')
            parts.append(f"\n{tool_num}. {tool_name} â€” {desc}")
            if example:
                parts.append(f"    {example}")
            tool_num += 1
    
    return "\n".join(parts)


# ============================================================
# TOOL ROUTER
# ============================================================

TOOLS = {
    "run_python": lambda args: tool_run_python(args.get("code", "")),
    "run_powershell": lambda args: tool_run_powershell(args.get("code", "")),
    "create_file": lambda args: tool_create_file(args.get("path", ""), args.get("content", "")),
    "read_file": lambda args: tool_read_file(args.get("path", "")),
    "list_files": lambda args: tool_list_files(args.get("path", ""), args.get("pattern", "*")),
    "organize_folder": lambda args: tool_organize_folder(args.get("path", "")),
    "find_duplicates": lambda args: tool_find_duplicates(args.get("path", "")),
    "disk_usage": lambda args: tool_disk_usage(args.get("path", "C:/")),
    "open_app": lambda args: tool_open_app(args.get("target", "")),
    "create_document": lambda args: tool_create_document(
        args.get("type", "word"), args.get("path", ""),
        args.get("title", ""), args.get("content", "")
    ),
    "web_search": lambda args: tool_web_search(args.get("query", "")),
    "parse_webpage": lambda args: tool_parse_webpage(args.get("url", ""), args.get("selector")),
    "stealth_browse": lambda args: tool_stealth_browse(
        args.get("url", ""), args.get("wait", 5), args.get("screenshot", False)
    ),
    "take_screenshot": lambda args: tool_take_screenshot(args.get("url", ""), args.get("path")),
    "browse_as_me": lambda args: tool_browse_as_me(
        args.get("url", ""), args.get("action", ""), args.get("wait", 5), args.get("screenshot", False)
    ),
    "system_info": lambda args: tool_system_info(),
    "clean_temp": lambda args: tool_clean_temp(args.get("dry_run", True)),
    "send_email": lambda args: tool_send_email(args.get("to", ""), args.get("subject", ""), args.get("body", "")),
    "n8n_list_workflows": lambda args: tool_n8n_list_workflows(args.get("query", ""), args.get("raw", False)),
    "n8n_get_workflow": lambda args: tool_n8n_get_workflow(args.get("id", ""), args.get("name", ""), args.get("raw", False)),
    "n8n_get_executions": lambda args: tool_n8n_get_executions(args.get("workflow_id", ""), args.get("limit", 5), args.get("raw", False)),
    "n8n_get_execution": lambda args: tool_n8n_get_execution(args.get("execution_id", ""), args.get("raw", False)),
    "n8n_run_workflow": lambda args: tool_n8n_run_workflow(args.get("workflow_id", ""), args.get("wait", True), args.get("raw", False)),
    "n8n_update_workflow": lambda args: tool_n8n_update_workflow(args.get("id", ""), args.get("workflow_json", {}), args.get("raw", False)),
    "n8n_validate_workflow": lambda args: tool_n8n_validate_workflow(args.get("workflow_json", {}), args.get("raw", False)),
    "n8n_create_workflow": lambda args: tool_n8n_create_workflow(
        args.get("name", "My Workflow"),
        args.get("description", ""),
        args.get("workflow_json"),
        args.get("raw", False),
    ),
    "n8n_activate_workflow": lambda args: tool_n8n_activate_workflow(args.get("id", ""), args.get("active", True)),
    "n8n_delete_workflow": lambda args: tool_n8n_delete_workflow(args.get("id", "")),
    "delete_files": lambda args: tool_delete_files(
        args.get("paths", []),
        args.get("permanent", False),
        args.get("confirm", False),
        args.get("allowed_folder")
    ),
    "clean_duplicates": lambda args: tool_clean_duplicates(
        args.get("path", ""),
        args.get("mode", "trash"),
        args.get("keep", "newest"),
        args.get("permanent", False),
        args.get("confirm", False),
        args.get("allowed_folder"),
    ),
    "move_files": lambda args: tool_move_files(args.get("files", [])),
    "restore_from_trash": lambda args: tool_restore_from_trash(args.get("path", "")),
    "list_trash": lambda args: tool_list_trash(args.get("drive")),
    "purge_trash": lambda args: tool_purge_trash(args.get("drive"), args.get("confirm", False)),
    "memory_add_fact": lambda args: tool_memory_add_fact(args.get("fact", "")),
    "memory_add_task": lambda args: tool_memory_add_task(args.get("task", "")),
    "memory_complete_task": lambda args: tool_memory_complete_task(args.get("task_index", -1)),
    "memory_show": lambda args: tool_memory_show(),
    "chat": lambda args: args.get("message", ""),
}

def init_tools():
    """Initialize tools: load skills and merge into TOOLS dict"""
    load_skills()
    # Merge skill tools into main TOOLS
    TOOLS.update(SKILL_TOOLS_MAP)

# ============================================================
# LLM COMMUNICATION
# ============================================================

def ask_ollama(user_message: str, history: list) -> str:
    """Send message to Ollama and get response.

    Routes through ProviderManager when available (supports local/api/auto).
    Falls back to direct Ollama HTTP call if ProviderManager is not initialised.
    """
    history.append({"role": "user", "content": user_message})

    # Build dynamic system prompt with memory context and code intelligence
    system_prompt = build_system_prompt(user_message)
    messages = [{"role": "system", "content": system_prompt}] + history[-20:]

    try:
        if _provider_mgr is not None:
            # --- NEW: route through ProviderManager ---
            llm_resp = _provider_mgr.chat(
                messages, temperature=0.3, max_tokens=4096, timeout=180
            )
            answer = llm_resp.content
            _dlog.log_provider(
                llm_resp.provider,
                fallback=_provider_mgr.used_fallback,
                ollama=_net_status.ollama,
                internet=_net_status.internet,
            )
        else:
            # --- LEGACY: direct Ollama call (backward compat) ---
            resp = requests.post(f"{OLLAMA_URL}/api/chat", json={
                "model": MODEL,
                "messages": messages,
                "stream": False,
                "options": {"temperature": 0.3, "num_predict": 4096}
            }, timeout=180)
            resp.encoding = 'utf-8'
            data = resp.json()
            answer = data.get("message", {}).get("content", "")
            if not answer:
                answer = data.get("response", "") or str(data)

        if not answer:
            answer = '{"tool": "chat", "args": {"message": "Empty response from LLM."}}'
        history.append({"role": "assistant", "content": answer})

        # Auto-save chat history after each exchange
        save_chat_history(history)

        return answer
    except requests.exceptions.ConnectionError:
        return '{"tool": "chat", "args": {"message": "Error: Cannot connect to Ollama. Is it running? Start Ollama first."}}'
    except requests.exceptions.Timeout:
        return '{"tool": "chat", "args": {"message": "Error: Ollama took too long to respond (180s timeout). Try a simpler request."}}'
    except ConnectionError as e:
        # ProviderManager raises ConnectionError when no provider is available
        return f'{{"tool": "chat", "args": {{"message": "No LLM available: {e}"}}}}'
    except NotImplementedError as e:
        # APIProvider raises NotImplementedError when api mode used but not yet implemented
        return f'{{"tool": "chat", "args": {{"message": "API provider not yet available: {e}"}}}}'
    except Exception as e:
        return f'{{"tool": "chat", "args": {{"message": "Error communicating with LLM: {e}"}}}}'


def extract_json(text: str) -> dict:
    """Extract JSON from model response"""
    text = text.strip()
    if "<think>" in text:
        idx = text.rfind("</think>")
        if idx != -1:
            text = text[idx + 8:].strip()
    if "```json" in text:
        text = text.split("```json")[1].split("```")[0].strip()
    elif "```" in text:
        parts = text.split("```")
        if len(parts) >= 3:
            text = parts[1].strip()
    start = text.find("{")
    end = text.rfind("}") + 1
    if start != -1 and end > start:
        try:
            return json.loads(text[start:end])
        except json.JSONDecodeError:
            json_str = text[start:end]
            json_str = json_str.replace("'", '"')
            try:
                return json.loads(json_str)
            except:
                pass
    return None


def _extract_path_from_message(msg: str) -> str:
    """Extract first Windows-style path from user message."""
    patterns = [
        r'([A-Za-z]:[/\\][^\s,"\'>]+)',
        r'([A-Za-z]:\\[^\s,"\'>]+)',
    ]
    for pat in patterns:
        match = re.search(pat, msg)
        if match:
            return match.group(1)
    return ""


def _controller_last_scanned_folder() -> str:
    """Read last_scanned_folder from controller policy if controller is present."""
    controller = globals().get("CONTROLLER")
    try:
        if controller and getattr(controller, "policy", None):
            return getattr(controller.policy, "last_scanned_folder", "") or ""
    except Exception:
        pass
    return ""


def _normalize_folder(path: str) -> str:
    if not path:
        return ""
    return os.path.normpath(os.path.abspath(path))


def _infer_allowed_folder(user_message: str, tool_args: dict) -> str:
    """Infer allowed folder for delete operations."""
    extracted = _extract_path_from_message(user_message)
    if extracted:
        candidate = extracted
        # If user message points to a file path, scope deletions to its parent folder.
        if os.path.isfile(candidate) or os.path.splitext(os.path.basename(candidate))[1]:
            candidate = os.path.dirname(candidate) or candidate
        return _normalize_folder(candidate)

    controller_folder = _controller_last_scanned_folder()
    if controller_folder:
        return _normalize_folder(controller_folder)

    if LAST_SCANNED_FOLDER:
        return _normalize_folder(LAST_SCANNED_FOLDER)

    return ""


# ============================================================
# LOGGING
# ============================================================

def log_interaction(user_msg: str, tool_name: str, result: str):
    """Log interactions to file"""
    log_file = os.path.join(LOGS_DIR, f"log_{datetime.datetime.now().strftime('%Y-%m-%d')}.jsonl")
    entry = {
        "timestamp": datetime.datetime.now().isoformat(),
        "user": user_msg,
        "tool": tool_name,
        "result": result[:500]
    }
    try:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(json.dumps(entry, ensure_ascii=False) + '\n')
    except:
        pass


# ============================================================
# PROCESS USER MESSAGE â€” AGENTIC LOOP (think â†’ act â†’ observe â†’ repeat)
# ============================================================

MAX_TOOL_STEPS = 5  # Maximum tool calls per user message (prevent infinite loops)

def process_message(user_message: str, history: list) -> dict:
    """
    Process a user message through an agentic loop.
    The agent can call multiple tools in sequence, feeding results back to Ollama.
    Returns dict with keys: response, tool_name, tool_result, raw, steps
    """
    steps = []  # Track all tool calls for this message
    last_tool_name = None
    last_tool_result = None

    # BrainLayer / Controller-first deterministic workflows (if available).
    # BrainLayer adds SLOW/FAST/CLARIFY routing on top of the controller.
    brain = globals().get("BRAIN")
    controller = globals().get("CONTROLLER")
    orchestrator = brain if brain is not None else controller
    if orchestrator is not None:
        try:
            if brain is not None:
                ctrl_result = brain.handle(user_message)
            else:
                ctrl_result = controller.handle_request(user_message)
            if ctrl_result.get("handled") or ctrl_result.get("awaiting_confirmation"):
                response_text = ctrl_result.get("response", "")
                tool_name = ctrl_result.get("tool_name", "brain" if brain is not None else "controller")
                log_interaction(user_message, tool_name, str(response_text))
                update_profile_from_conversation(user_message, tool_name, str(response_text))
                # Serialize safely: convert non-JSON-serializable objects (e.g. PlanStep) to str.
                try:
                    raw_str = json.dumps(ctrl_result, ensure_ascii=False, default=str)
                except Exception:
                    raw_str = str(ctrl_result)
                return {
                    "response": response_text,
                    "tool_name": ctrl_result.get("tool_name"),
                    "tool_result": ctrl_result.get("tool_result"),
                    "raw": raw_str,
                    "thinking": ctrl_result.get("thinking", ""),
                    "steps": ctrl_result.get("steps", []),
                }
        except Exception as e:
            import traceback
            log_interaction(user_message, "brain_error" if brain else "controller_error",
                            str(e) + "\n" + traceback.format_exc()[:500])
    
    for step in range(MAX_TOOL_STEPS):
        answer = ask_ollama(user_message if step == 0 else f"[Tool result for {last_tool_name}]:\n{last_tool_result[:3000]}\n\nBased on this result, what should I do next? If the task is complete, use chat tool to tell Jane the result in her language.", history)
        parsed = extract_json(answer)

        if not parsed:
            # Try fallback guessing
            text = answer.strip()
            if "<think>" in text:
                idx = text.rfind("</think>")
                if idx != -1:
                    text = text[idx + 8:].strip()
            
            # If model returned plain text without JSON, treat it as a chat response
            if text and not text.startswith("{"):
                log_interaction(user_message, "raw", text)
                update_profile_from_conversation(user_message, "raw", text)
                return {"response": text, "tool_name": None, "tool_result": None, "raw": answer, "steps": steps}
            
            lower = text.lower()
            ps_prefixes = ["start-process", "start ", "get-", "set-", "new-", "remove-",
                           "invoke-", "mkdir ", "dir ", "copy ", "move ", "del "]
            for prefix in ps_prefixes:
                if lower.startswith(prefix):
                    parsed = {"tool": "run_powershell", "args": {"code": text}}
                    break
            if not parsed:
                py_prefixes = ["import ", "from ", "print(", "open(", "os.", "with "]
                for prefix in py_prefixes:
                    if lower.startswith(prefix):
                        parsed = {"tool": "run_python", "args": {"code": text}}
                        break

        if not parsed:
            # Could not parse anything â€” return raw answer
            log_interaction(user_message, "raw", answer)
            return {"response": answer, "tool_name": None, "tool_result": None, "raw": answer, "steps": steps}

        tool_name = parsed.get("tool", "")
        tool_args = parsed.get("args", {})
        thinking = parsed.get("thinking", "")

        # ============================================================
        # TOOL ALIAS / FUZZY MATCHING
        # ============================================================
        tool_name = _resolve_tool_alias(tool_name)

        # --- CHAT: direct response to user ---
        if tool_name == "chat":
            message = tool_args.get("message", answer)
            # Prepend thinking if available (for transparency in web UI)
            log_interaction(user_message, "chat", message)
            update_profile_from_conversation(user_message, "chat", message)
            return {"response": message, "tool_name": "chat", "tool_result": None, "raw": answer, 
                    "thinking": thinking, "steps": steps}

        # --- EXECUTE TOOL ---
        elif tool_name in TOOLS:
            if tool_name in {"delete_files", "clean_duplicates"} and not tool_args.get("allowed_folder"):
                inferred_folder = _infer_allowed_folder(user_message, tool_args)
                if inferred_folder:
                    tool_args["allowed_folder"] = inferred_folder
                else:
                    ask_folder_msg = (
                        "I need the target folder before deleting. "
                        "Please provide the folder path (for example: E:/Alexander)."
                    )
                    log_interaction(user_message, "chat", ask_folder_msg)
                    update_profile_from_conversation(user_message, "chat", ask_folder_msg)
                    return {
                        "response": ask_folder_msg,
                        "tool_name": "chat",
                        "tool_result": None,
                        "raw": answer,
                        "thinking": thinking,
                        "steps": steps,
                    }

            try:
                result = TOOLS[tool_name](tool_args)
            except Exception as e:
                result = f"Error executing {tool_name}: {e}"
            
            steps.append({"tool": tool_name, "args": tool_args, "result": result[:1000], "thinking": thinking})
            log_interaction(user_message, tool_name, result)
            update_profile_from_conversation(user_message, tool_name, result)
            save_chat_history(history)
            
            last_tool_name = tool_name
            last_tool_result = result
            
            # Continue the loop â€” let agent decide if more steps are needed
            # (unless this is the last allowed step)
            if step < MAX_TOOL_STEPS - 1:
                continue
            else:
                # Max steps reached â€” return what we have
                return {"response": None, "tool_name": tool_name, "tool_result": result, "raw": answer,
                        "thinking": thinking, "steps": steps}

        else:
            # Tool not found even after alias resolution â€” ask agent to retry
            log_interaction(user_message, "unknown", answer)
            if step < MAX_TOOL_STEPS - 1:
                # Feed error back and let agent try again
                history.append({"role": "user", "content": f"Error: tool '{tool_name}' does not exist. Available tools: {', '.join(sorted(TOOLS.keys()))}. Please try again with a valid tool name."})
                last_tool_name = "error"
                last_tool_result = f"Tool '{tool_name}' not found"
                continue
            else:
                return {
                    "response": f"I couldn't find the right tool for this request. Let me try differently â€” please rephrase.",
                    "tool_name": None, "tool_result": None, "raw": answer, "steps": steps
                }

    # If we exit the loop (all steps used on tools, last step was a tool call)
    # Ask Ollama for a final summary
    if last_tool_result:
        summary_answer = ask_ollama(
            f"[Tool result for {last_tool_name}]:\n{last_tool_result[:3000]}\n\nThe task is done. Summarize what happened for Jane in her language using chat tool.",
            history
        )
        summary_parsed = extract_json(summary_answer)
        if summary_parsed and summary_parsed.get("tool") == "chat":
            return {"response": summary_parsed["args"].get("message", last_tool_result), 
                    "tool_name": last_tool_name, "tool_result": last_tool_result, "raw": summary_answer, "steps": steps}
        return {"response": None, "tool_name": last_tool_name, "tool_result": last_tool_result, 
                "raw": summary_answer, "steps": steps}

    return {"response": "I couldn't process this request. Please try rephrasing.", 
            "tool_name": None, "tool_result": None, "raw": "", "steps": steps}


# Tool alias map â€” moved outside process_message for clarity
TOOL_ALIASES = {
    "browse_url": "stealth_browse", "open_url": "stealth_browse", "browse_website": "stealth_browse",
    "open_website": "stealth_browse", "visit_url": "stealth_browse", "goto_url": "stealth_browse",
    "navigate": "stealth_browse", "browse": "stealth_browse", "open_browser": "stealth_browse",
    "open_page": "stealth_browse", "load_page": "stealth_browse",
    "fetch_url": "parse_webpage", "fetch_page": "parse_webpage", "scrape_page": "parse_webpage",
    "scrape_url": "parse_webpage", "get_page": "parse_webpage", "download_page": "parse_webpage",
    "create_text_file": "create_file", "write_file": "create_file", "save_file": "create_file",
    "make_file": "create_file", "new_file": "create_file", "touch": "create_file",
    "open_file": "read_file", "cat_file": "read_file", "view_file": "read_file", "show_file": "read_file",
    "delete_file": "delete_files", "remove_file": "delete_files", "remove_files": "delete_files", "rm": "delete_files",
    "rename_file": "move_files", "rename_files": "move_files", "move_file": "move_files", "mv": "move_files",
    "ls": "list_files", "dir": "list_files", "list_directory": "list_files", "show_files": "list_files",
    "google": "web_search", "search": "web_search", "search_web": "web_search", "google_search": "web_search",
    "screenshot": "take_screenshot", "capture_screen": "take_screenshot",
    "sysinfo": "system_info", "sys_info": "system_info", "pc_info": "system_info",
    "email": "send_email", "mail": "send_email", "send_mail": "send_email", "gmail": "send_email",
    "python": "run_python", "exec_python": "run_python", "powershell": "run_powershell",
    "shell": "run_powershell", "cmd": "run_powershell", "terminal": "run_powershell",
    "create_doc": "create_document", "create_word": "create_document", "create_pdf": "create_document",
    "remember": "memory_add_fact", "add_task": "memory_add_task", "show_memory": "memory_show",
    "send_telegram": "telegram_send", "tg_send": "telegram_send", "send_message": "telegram_send",
}

def _resolve_tool_alias(tool_name: str) -> str:
    """Resolve tool aliases and fuzzy matches"""
    if tool_name in TOOLS:
        return tool_name
    if tool_name in TOOL_ALIASES:
        real = TOOL_ALIASES[tool_name]
        if real in TOOLS:
            return real
    # Fuzzy match
    if tool_name != "chat":
        tool_lower = tool_name.lower().replace("_", "").replace("-", "")
        for real_name in TOOLS:
            real_lower = real_name.lower().replace("_", "").replace("-", "")
            if tool_lower == real_lower or tool_lower in real_lower or real_lower in tool_lower:
                return real_name
    return tool_name


# ============================================================
# WEB UI - HTML/CSS/JS
# ============================================================

HTML_PAGE = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Jane's AI Agent v4</title>
<style>
:root {
    --bg-primary: #0d1117;
    --bg-secondary: #161b22;
    --bg-tertiary: #1c2128;
    --bg-input: #21262d;
    --border: #30363d;
    --text-primary: #e6edf3;
    --text-secondary: #8b949e;
    --text-muted: #6e7681;
    --accent: #58a6ff;
    --accent-hover: #79c0ff;
    --accent-bg: rgba(56,139,253,0.15);
    --green: #3fb950;
    --green-bg: rgba(46,160,67,0.15);
    --orange: #d29922;
    --orange-bg: rgba(187,128,9,0.15);
    --red: #f85149;
    --purple: #bc8cff;
    --purple-bg: rgba(188,140,255,0.15);
    --shadow: 0 8px 24px rgba(0,0,0,0.4);
    --radius: 12px;
    --radius-sm: 8px;
}

* { margin: 0; padding: 0; box-sizing: border-box; }

body {
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', 'Noto Sans', Helvetica, Arial, sans-serif;
    background: var(--bg-primary);
    color: var(--text-primary);
    height: 100vh;
    display: flex;
    flex-direction: column;
    overflow: hidden;
}

/* Header */
.header {
    background: var(--bg-secondary);
    border-bottom: 1px solid var(--border);
    padding: 12px 24px;
    display: flex;
    align-items: center;
    justify-content: space-between;
    flex-shrink: 0;
    z-index: 10;
}
.header-left {
    display: flex;
    align-items: center;
    gap: 12px;
}
.logo {
    width: 36px;
    height: 36px;
    background: linear-gradient(135deg, var(--accent), var(--purple));
    border-radius: 10px;
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 18px;
    font-weight: 700;
    color: white;
}
.header h1 {
    font-size: 16px;
    font-weight: 600;
    color: var(--text-primary);
}
.header .subtitle {
    font-size: 12px;
    color: var(--text-secondary);
}
.status-badge {
    display: flex;
    align-items: center;
    gap: 6px;
    padding: 4px 12px;
    border-radius: 20px;
    font-size: 12px;
    font-weight: 500;
}
.status-badge.connected {
    background: var(--green-bg);
    color: var(--green);
}
.status-badge.disconnected {
    background: rgba(248,81,73,0.15);
    color: var(--red);
}
.status-dot {
    width: 8px;
    height: 8px;
    border-radius: 50%;
    background: currentColor;
    animation: pulse 2s ease-in-out infinite;
}
@keyframes pulse {
    0%, 100% { opacity: 1; }
    50% { opacity: 0.4; }
}

/* Chat Container */
.chat-container {
    flex: 1;
    overflow-y: auto;
    padding: 20px 0;
    scroll-behavior: smooth;
}
.chat-container::-webkit-scrollbar {
    width: 6px;
}
.chat-container::-webkit-scrollbar-track {
    background: transparent;
}
.chat-container::-webkit-scrollbar-thumb {
    background: var(--border);
    border-radius: 3px;
}

/* Messages */
.message-group {
    max-width: 820px;
    margin: 0 auto;
    padding: 0 24px;
}
.message {
    padding: 16px 0;
    display: flex;
    gap: 16px;
    animation: fadeIn 0.3s ease;
}
@keyframes fadeIn {
    from { opacity: 0; transform: translateY(8px); }
    to { opacity: 1; transform: translateY(0); }
}
.message-avatar {
    width: 32px;
    height: 32px;
    border-radius: 8px;
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 14px;
    font-weight: 600;
    flex-shrink: 0;
    margin-top: 2px;
}
.message.user .message-avatar {
    background: linear-gradient(135deg, #6e40c9, #8957e5);
    color: white;
}
.message.assistant .message-avatar {
    background: linear-gradient(135deg, var(--accent), #1f6feb);
    color: white;
}
.message-content {
    flex: 1;
    min-width: 0;
}
.message-content .name {
    font-size: 13px;
    font-weight: 600;
    color: var(--text-secondary);
    margin-bottom: 6px;
}
.message-content .text {
    font-size: 14px;
    line-height: 1.65;
    color: var(--text-primary);
    word-wrap: break-word;
    white-space: pre-wrap;
}

/* Tool blocks */
.tool-block {
    margin: 8px 0;
    border-radius: var(--radius-sm);
    border: 1px solid var(--border);
    overflow: hidden;
}
.tool-header {
    padding: 8px 14px;
    background: var(--bg-tertiary);
    display: flex;
    align-items: center;
    gap: 8px;
    font-size: 12px;
    font-weight: 600;
    color: var(--text-secondary);
    border-bottom: 1px solid var(--border);
}
.tool-icon {
    font-size: 14px;
}
.tool-body {
    padding: 12px 14px;
    background: var(--bg-secondary);
    font-family: 'JetBrains Mono', 'Cascadia Code', 'Fira Code', monospace;
    font-size: 12.5px;
    line-height: 1.55;
    color: var(--text-primary);
    white-space: pre-wrap;
    word-break: break-word;
    max-height: 400px;
    overflow-y: auto;
}
.tool-body::-webkit-scrollbar {
    width: 4px;
}
.tool-body::-webkit-scrollbar-thumb {
    background: var(--border);
    border-radius: 2px;
}

/* Loading */
.typing-indicator {
    display: flex;
    align-items: center;
    gap: 4px;
    padding: 8px 0;
}
.typing-dot {
    width: 6px;
    height: 6px;
    border-radius: 50%;
    background: var(--text-muted);
    animation: typingBounce 1.4s ease-in-out infinite;
}
.typing-dot:nth-child(2) { animation-delay: 0.2s; }
.typing-dot:nth-child(3) { animation-delay: 0.4s; }
@keyframes typingBounce {
    0%, 60%, 100% { transform: translateY(0); }
    30% { transform: translateY(-6px); }
}

/* Input Area */
.input-area {
    background: var(--bg-secondary);
    border-top: 1px solid var(--border);
    padding: 16px 24px 20px;
    flex-shrink: 0;
}
.input-wrapper {
    max-width: 820px;
    margin: 0 auto;
    position: relative;
}
.input-box {
    width: 100%;
    background: var(--bg-input);
    border: 1px solid var(--border);
    border-radius: var(--radius);
    padding: 14px 56px 14px 18px;
    color: var(--text-primary);
    font-size: 14px;
    font-family: inherit;
    line-height: 1.5;
    resize: none;
    outline: none;
    min-height: 52px;
    max-height: 200px;
    transition: border-color 0.2s, box-shadow 0.2s;
}
.input-box:focus {
    border-color: var(--accent);
    box-shadow: 0 0 0 3px var(--accent-bg);
}
.input-box::placeholder {
    color: var(--text-muted);
}
.send-btn {
    position: absolute;
    right: 10px;
    bottom: 10px;
    width: 36px;
    height: 36px;
    border-radius: 8px;
    border: none;
    background: var(--accent);
    color: white;
    cursor: pointer;
    display: flex;
    align-items: center;
    justify-content: center;
    transition: background 0.2s, transform 0.1s;
}
.send-btn:hover { background: var(--accent-hover); }
.send-btn:active { transform: scale(0.95); }
.send-btn:disabled {
    background: var(--bg-tertiary);
    color: var(--text-muted);
    cursor: not-allowed;
}
.send-btn svg {
    width: 18px;
    height: 18px;
}

/* Welcome */
.welcome {
    max-width: 820px;
    margin: 0 auto;
    padding: 60px 24px 20px;
    text-align: center;
}
.welcome h2 {
    font-size: 24px;
    font-weight: 600;
    margin-bottom: 8px;
    background: linear-gradient(135deg, var(--accent), var(--purple));
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
}
.welcome p {
    color: var(--text-secondary);
    font-size: 14px;
    margin-bottom: 32px;
}
.quick-actions {
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
    gap: 10px;
    max-width: 640px;
    margin: 0 auto;
}
.quick-action {
    background: var(--bg-secondary);
    border: 1px solid var(--border);
    border-radius: var(--radius-sm);
    padding: 14px 16px;
    text-align: left;
    cursor: pointer;
    transition: border-color 0.2s, background 0.2s;
}
.quick-action:hover {
    border-color: var(--accent);
    background: var(--accent-bg);
}
.quick-action .qa-icon {
    font-size: 18px;
    margin-bottom: 6px;
}
.quick-action .qa-title {
    font-size: 13px;
    font-weight: 600;
    color: var(--text-primary);
    margin-bottom: 2px;
}
.quick-action .qa-desc {
    font-size: 12px;
    color: var(--text-secondary);
}

/* Footer info */
.footer-info {
    text-align: center;
    font-size: 11px;
    color: var(--text-muted);
    padding: 4px 0 0;
}
</style>
</head>
<body>

<div class="header">
    <div class="header-left">
        <div class="logo">J</div>
        <div>
            <h1>Jane's AI Agent</h1>
            <div class="subtitle">v4.1 &middot; Local &middot; Skills &middot; Telegram &middot; Chrome Profile &middot; CRM</div>
        </div>
    </div>
    <div style="display:flex;align-items:center;gap:10px;">
        <button onclick="clearHistory()" style="background:var(--bg-tertiary);border:1px solid var(--border);color:var(--text-secondary);padding:4px 10px;border-radius:6px;cursor:pointer;font-size:11px;" title="Clear chat history">Clear</button>
        <button onclick="sendQuick('Show my memory and profile')" style="background:var(--bg-tertiary);border:1px solid var(--border);color:var(--text-secondary);padding:4px 10px;border-radius:6px;cursor:pointer;font-size:11px;" title="Show memory">Memory</button>
        <div id="status" class="status-badge disconnected">
            <span class="status-dot"></span>
            <span id="statusText">Connecting...</span>
        </div>
    </div>
</div>

<div class="chat-container" id="chat">
    <div class="welcome" id="welcome">
        <h2>Ready to help, Jane</h2>
        <p>Write in any language, or send tasks from Telegram. I'll respond in yours.</p>
        <div class="quick-actions">
            <div class="quick-action" onclick="sendQuick('Show system info: CPU, RAM, disk')">
                <div class="qa-icon">&#128187;</div>
                <div class="qa-title">System Info</div>
                <div class="qa-desc">CPU, RAM, disk usage</div>
            </div>
            <div class="quick-action" onclick="sendQuick('Show my n8n workflows')">
                <div class="qa-icon">&#9889;</div>
                <div class="qa-title">n8n Workflows</div>
                <div class="qa-desc">List all automations</div>
            </div>
            <div class="quick-action" onclick="sendQuick('Show Telegram bot status')">
                <div class="qa-icon">&#128172;</div>
                <div class="qa-title">Telegram Status</div>
                <div class="qa-desc">Bot connection &amp; messages</div>
            </div>
            <div class="quick-action" onclick="sendQuick('Search Swedish municipalities in Stockholm region')">
                <div class="qa-icon">&#127474;&#127480;</div>
                <div class="qa-title">Kommun Parser</div>
                <div class="qa-desc">Swedish municipality data</div>
            </div>
            <div class="quick-action" onclick="sendQuick('Show CRM statistics')">
                <div class="qa-icon">&#128188;</div>
                <div class="qa-title">BoostCamp CRM</div>
                <div class="qa-desc">Startups &amp; investors</div>
            </div>
            <div class="quick-action" onclick="sendQuick('Show all available n8n workflow templates')">
                <div class="qa-icon">&#128203;</div>
                <div class="qa-title">n8n Templates</div>
                <div class="qa-desc">Ready-made automations</div>
            </div>
            <div class="quick-action" onclick="sendQuick('Show today calendar events')">
                <div class="qa-icon">&#128197;</div>
                <div class="qa-title">Calendar Today</div>
                <div class="qa-desc">Google Calendar events</div>
            </div>
            <div class="quick-action" onclick="sendQuick('Check voice input status')">
                <div class="qa-icon">&#127908;</div>
                <div class="qa-title">Voice Input</div>
                <div class="qa-desc">Whisper speech-to-text</div>
            </div>
        </div>
    </div>
    <div class="message-group" id="messages"></div>
</div>

<div class="input-area">
    <div class="input-wrapper">
        <textarea class="input-box" id="input" placeholder="Write anything... (Enter to send, Shift+Enter for new line)" rows="1"></textarea>
        <button class="send-btn" id="sendBtn" onclick="sendMessage()">
            <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round">
                <line x1="22" y1="2" x2="11" y2="13"></line>
                <polygon points="22 2 15 22 11 13 2 9 22 2"></polygon>
            </svg>
        </button>
    </div>
    <div class="footer-info">Ollama &middot; local model &middot; skills &middot; Telegram &middot; CRM &middot; your data stays on your PC</div>
</div>

<script>
const chat = document.getElementById('chat');
const messages = document.getElementById('messages');
const input = document.getElementById('input');
const sendBtn = document.getElementById('sendBtn');
const welcome = document.getElementById('welcome');
const statusEl = document.getElementById('status');
const statusText = document.getElementById('statusText');

let isWaiting = false;

// Auto-resize textarea
input.addEventListener('input', () => {
    input.style.height = 'auto';
    input.style.height = Math.min(input.scrollHeight, 200) + 'px';
});

// Enter to send
input.addEventListener('keydown', (e) => {
    if (e.key === 'Enter' && !e.shiftKey) {
        e.preventDefault();
        sendMessage();
    }
});

// Check Ollama status
async function checkStatus() {
    try {
        const r = await fetch('/api/status');
        const data = await r.json();
        if (data.ollama) {
            statusEl.className = 'status-badge connected';
            statusText.textContent = data.model || 'Connected';
        } else {
            statusEl.className = 'status-badge disconnected';
            statusText.textContent = 'Ollama offline';
        }
    } catch {
        statusEl.className = 'status-badge disconnected';
        statusText.textContent = 'Server error';
    }
}
checkStatus();
setInterval(checkStatus, 30000);

function escapeHtml(text) {
    const div = document.createElement('div');
    div.textContent = text;
    return div.innerHTML;
}

function addMessage(role, content, toolName, toolResult, thinking) {
    if (welcome) welcome.style.display = 'none';

    const msg = document.createElement('div');
    msg.className = `message ${role}`;

    const avatar = document.createElement('div');
    avatar.className = 'message-avatar';
    avatar.textContent = role === 'user' ? 'J' : 'AI';

    const contentDiv = document.createElement('div');
    contentDiv.className = 'message-content';

    const nameDiv = document.createElement('div');
    nameDiv.className = 'name';
    nameDiv.textContent = role === 'user' ? 'Jane' : 'Agent';
    contentDiv.appendChild(nameDiv);

    // Show thinking process (if available)
    if (thinking) {
        const thinkDiv = document.createElement('div');
        thinkDiv.style.cssText = 'color: var(--text-muted); font-size: 12px; font-style: italic; margin-bottom: 6px; padding: 4px 8px; border-left: 2px solid var(--border); opacity: 0.8;';
        thinkDiv.textContent = thinking;
        contentDiv.appendChild(thinkDiv);
    }

    if (content) {
        const textDiv = document.createElement('div');
        textDiv.className = 'text';
        textDiv.textContent = content;
        contentDiv.appendChild(textDiv);
    }

    if (toolName && toolName !== 'chat') {
        const toolBlock = document.createElement('div');
        toolBlock.className = 'tool-block';

        const toolIcons = {
            'run_python': '\u{1F40D}',
            'run_powershell': '\u{1F4BB}',
            'create_file': '\u{1F4DD}',
            'read_file': '\u{1F4D6}',
            'list_files': '\u{1F4C2}',
            'organize_folder': '\u{1F4E6}',
            'find_duplicates': '\u{1F50D}',
            'disk_usage': '\u{1F4CA}',
            'open_app': '\u{1F680}',
            'create_document': '\u{1F4C4}',
            'web_search': '\u{1F310}',
            'parse_webpage': '\u{1F578}',
            'stealth_browse': '\u{1F575}',
            'take_screenshot': '\u{1F4F7}',
            'system_info': '\u{2699}',
            'clean_temp': '\u{1F9F9}',
            'send_email': '\u{2709}',
            'n8n_list_workflows': '\u{26A1}',
            'n8n_get_workflow': '\u{26A1}',
            'n8n_get_executions': '\u{1F4DD}',
            'n8n_get_execution': '\u{1F50E}',
            'n8n_run_workflow': '\u{25B6}',
            'n8n_update_workflow': '\u{1F527}',
            'n8n_validate_workflow': '\u{2705}',
            'n8n_create_workflow': '\u{1F916}',
            'n8n_activate_workflow': '\u{2705}',
            'n8n_delete_workflow': '\u{274C}',
            'delete_files': '\u{1F5D1}',
            'move_files': '\u{1F4E4}',
            'memory_add_fact': '\u{1F4BE}',
            'memory_add_task': '\u{1F4CB}',
            'memory_complete_task': '\u{2705}',
            'memory_show': '\u{1F9E0}',
        };

        const icon = toolIcons[toolName] || '\u{1F527}';

        toolBlock.innerHTML = `
            <div class="tool-header">
                <span class="tool-icon">${icon}</span>
                <span>${escapeHtml(toolName)}</span>
            </div>
            <div class="tool-body">${escapeHtml(toolResult || '')}</div>
        `;
        contentDiv.appendChild(toolBlock);
    }

    msg.appendChild(avatar);
    msg.appendChild(contentDiv);
    messages.appendChild(msg);

    chat.scrollTop = chat.scrollHeight;
}

function addTyping() {
    const msg = document.createElement('div');
    msg.className = 'message assistant';
    msg.id = 'typing';

    msg.innerHTML = `
        <div class="message-avatar">AI</div>
        <div class="message-content">
            <div class="name">Agent</div>
            <div class="typing-indicator">
                <div class="typing-dot"></div>
                <div class="typing-dot"></div>
                <div class="typing-dot"></div>
            </div>
        </div>
    `;
    messages.appendChild(msg);
    chat.scrollTop = chat.scrollHeight;
}

function removeTyping() {
    const el = document.getElementById('typing');
    if (el) el.remove();
}

function sendQuick(text) {
    input.value = text;
    sendMessage();
}

async function sendMessage() {
    const text = input.value.trim();
    if (!text || isWaiting) return;

    isWaiting = true;
    sendBtn.disabled = true;
    input.value = '';
    input.style.height = 'auto';

    addMessage('user', text);
    addTyping();

    try {
        const resp = await fetch('/api/chat', {
            method: 'POST',
            headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({message: text})
        });
        const data = await resp.json();

        removeTyping();

        // Show intermediate steps if any (multi-step reasoning)
        if (data.steps && data.steps.length > 0) {
            for (const step of data.steps) {
                // Don't re-show the last step if it matches the final result
                if (step === data.steps[data.steps.length - 1] && !data.response) {
                    addMessage('assistant', step.thinking || null, step.tool, step.result);
                } else {
                    addMessage('assistant', step.thinking || null, step.tool, step.result);
                }
            }
            // Show final response if different from last step
            if (data.response) {
                addMessage('assistant', data.response);
            }
        } else if (data.response) {
            addMessage('assistant', data.response, data.tool_name, data.tool_result, data.thinking);
        } else if (data.tool_result) {
            addMessage('assistant', null, data.tool_name, data.tool_result, data.thinking);
        } else {
            addMessage('assistant', data.raw || 'No response');
        }
    } catch (err) {
        removeTyping();
        addMessage('assistant', 'Connection error: ' + err.message);
    }

    isWaiting = false;
    sendBtn.disabled = false;
    input.focus();
}

async function clearHistory() {
    if (!confirm('Clear chat history? Memory and profile will be kept.')) return;
    await fetch('/api/clear', {method: 'POST'});
    messages.innerHTML = '';
    if (welcome) welcome.style.display = '';
}

// Focus on load
input.focus();
</script>
</body>
</html>"""


# ============================================================
# WEB SERVER
# ============================================================

class AgentHandler(BaseHTTPRequestHandler):
    """HTTP request handler for the agent web UI"""

    def log_message(self, format, *args):
        """Suppress default HTTP logging (too noisy)"""
        pass

    def _send_json(self, data, code=200):
        self.send_response(code)
        self.send_header('Content-Type', 'application/json; charset=utf-8')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(json.dumps(data, ensure_ascii=False).encode('utf-8'))

    def _send_html(self, html, code=200):
        self.send_response(code)
        self.send_header('Content-Type', 'text/html; charset=utf-8')
        self.end_headers()
        self.wfile.write(html.encode('utf-8'))

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'GET, POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()

    def _serve_file(self, filename, content_type, binary=False):
        filepath = os.path.join(FRONTEND_DIR, filename)
        if not os.path.exists(filepath):
            self.send_error(404)
            return
        try:
            mode = 'rb' if binary else 'r'
            encoding = None if binary else 'utf-8'
            with open(filepath, mode, encoding=encoding) as f:
                content = f.read()
            self.send_response(200)
            self.send_header('Content-Type', f'{content_type}; charset=utf-8')
            self.send_header('Access-Control-Allow-Origin', '*')
            self.end_headers()
            if binary:
                self.wfile.write(content)
            else:
                self.wfile.write(content.encode('utf-8'))
        except Exception as e:
            self.send_error(500, str(e))

    def do_GET(self):
        parsed = urlparse(self.path)
        path = parsed.path

        # Serve static UI files if FRONTEND_DIR exists
        if os.path.isdir(FRONTEND_DIR):
            if path == '/' or path == '':
                self._serve_file('index.html', 'text/html')
                return
            elif path.endswith('.css'):
                self._serve_file(path.lstrip('/'), 'text/css')
                return
            elif path.endswith('.js') and not path.startswith('/api'):
                self._serve_file(path.lstrip('/'), 'application/javascript')
                return
        else:
            if path == '/' or path == '':
                self._send_html(HTML_PAGE)
                return

        if parsed.path == '/' or parsed.path == '':
            self._send_html(HTML_PAGE)

        elif parsed.path == '/api/status':
            # Check Ollama connection + new architecture status
            _net_status.refresh()
            try:
                resp = requests.get(f"{OLLAMA_URL}/api/tags", timeout=3)
                models = [m["name"] for m in resp.json().get("models", [])]
                has_model = MODEL in models or f"{MODEL}:latest" in models
                self._send_json({
                    "ollama": True,
                    "model": MODEL if has_model else None,
                    "models": models,
                    "has_model": has_model,
                    "tools": len(TOOLS),
                    "skills": len(LOADED_SKILLS),
                    "skill_names": list(LOADED_SKILLS.keys()),
                    "version": VERSION,
                    # New architecture fields
                    "provider_mode": _agent_cfg.provider_mode,
                    "knowledge_mode": _agent_cfg.knowledge_mode,
                    "internet": _net_status.internet,
                    "effective_mode": _net_status.effective_mode,
                })
            except:
                self._send_json({
                    "ollama": False, "model": None,
                    "tools": len(TOOLS), "skills": len(LOADED_SKILLS),
                    "version": VERSION,
                    "provider_mode": _agent_cfg.provider_mode,
                    "knowledge_mode": _agent_cfg.knowledge_mode,
                    "internet": _net_status.internet,
                    "effective_mode": _net_status.effective_mode,
                })

        elif parsed.path == '/api/tools':
            tools_list = sorted(TOOLS.keys())
            self._send_json({"tools": tools_list, "count": len(tools_list)})

        elif parsed.path == '/api/history':
            with conversation_lock:
                self._send_json({"history": conversation_history[-50:]})

        elif parsed.path == '/api/skills':
            skills_info = {}
            for name, module in LOADED_SKILLS.items():
                skills_info[name] = {
                    "version": getattr(module, 'SKILL_VERSION', '?'),
                    "description": getattr(module, 'SKILL_DESCRIPTION', ''),
                    "tools": list(getattr(module, 'SKILL_TOOLS', {}).keys())
                }
            self._send_json({"skills": skills_info, "count": len(skills_info)})

        elif parsed.path == '/api/memory':
            profile = load_user_profile()
            tasks = load_tasks()
            self._send_json({
                "profile": profile,
                "tasks": tasks,
                "history_length": len(conversation_history)
            })

        else:
            self.send_error(404)

    def do_POST(self):
        parsed = urlparse(self.path)

        if parsed.path == '/api/chat':
            content_length = int(self.headers.get('Content-Length', 0))
            body = self.rfile.read(content_length).decode('utf-8')

            try:
                data = json.loads(body)
            except:
                self._send_json({"error": "Invalid JSON"}, 400)
                return

            user_message = data.get("message", "").strip()
            if not user_message:
                self._send_json({"error": "Empty message"}, 400)
                return

            with conversation_lock:
                result = process_message(user_message, conversation_history)

            self._send_json(result)

        elif parsed.path == '/api/clear':
            with conversation_lock:
                conversation_history.clear()
                save_chat_history([])
            self._send_json({"status": "cleared"})

        else:
            self.send_error(404)


# ============================================================
# MAIN
# ============================================================

def main():
    import sys
    
    # Parse command-line arguments
    background_mode = '--background' in sys.argv or '--bg' in sys.argv or '--tray' in sys.argv
    no_browser = '--no-browser' in sys.argv or background_mode
    telegram_only = '--telegram-only' in sys.argv
    
    # Set low process priority in background mode (Windows)
    if background_mode:
        try:
            import ctypes
            # BELOW_NORMAL_PRIORITY_CLASS = 0x00004000
            ctypes.windll.kernel32.SetPriorityClass(
                ctypes.windll.kernel32.GetCurrentProcess(), 0x00004000
            )
        except:
            pass  # Not Windows or no ctypes
    
    # Initialize tools and skills FIRST
    init_tools()

    # Initialize controller layer (optional hardening path).
    global CONTROLLER, BRAIN
    try:
        from controller import create_controller
        CONTROLLER = create_controller(MEMORY_DIR, TOOLS)
        print("  Controller: enabled")
    except Exception as e:
        CONTROLLER = None
        print(f"  Controller: disabled ({e})")

    # Initialize BrainLayer (SLOW/FAST/CLARIFY routing) on top of controller.
    if CONTROLLER is not None:
        try:
            from brain.brain_layer import BrainLayer
            BRAIN = BrainLayer(
                CONTROLLER,
                require_confirmation=False,
                memory_dir=Path(MEMORY_DIR),
                provider=_provider_mgr._local if _provider_mgr is not None else None,
            )
            print("  BrainLayer: enabled")
        except Exception as e:
            BRAIN = None
            print(f"  BrainLayer: disabled ({e})")

    # Initialize CodeIntelligence (dependency graph + smart context + change tracking)
    global CODE_INTELLIGENCE
    try:
        from brain.code_intelligence import CodeIntelligence
        CODE_INTELLIGENCE = CodeIntelligence(AGENT_DIR)
        indexed = CODE_INTELLIGENCE.index_all()
        CODE_INTELLIGENCE.start_watching()
        print(f"  CodeIntelligence: indexed {indexed} symbols, watching for changes")
    except Exception as e:
        CODE_INTELLIGENCE = None
        print(f"  CodeIntelligence: disabled ({e})")

    mode_str = "BACKGROUND" if background_mode else "NORMAL"
    if telegram_only:
        mode_str = "TELEGRAM-ONLY"
    
    print(f"""
 ================================================================
     Jane's AI Agent v{VERSION} - Web UI + Telegram + Skills
     Mode: {mode_str}
     Model: {MODEL}
     Built-in tools: 30 | Skill tools: {len(SKILL_TOOLS_MAP)}
     Total tools: {len(TOOLS)} available
     Skills loaded: {len(LOADED_SKILLS)} ({', '.join(LOADED_SKILLS.keys()) if LOADED_SKILLS else 'none'})
     URL:   http://localhost:{WEB_PORT}
 ================================================================
""")

    if background_mode:
        print("  Mode: BACKGROUND (low CPU priority, no browser, Telegram active)")
        print("  The child can play games â€” agent works silently in background.")
        print("  Send tasks from Telegram, results come back to Telegram.\n")

    # Check Ollama
    try:
        resp = requests.get(f"{OLLAMA_URL}/api/tags", timeout=5)
        models = [m["name"] for m in resp.json().get("models", [])]
        if MODEL not in models and f"{MODEL}:latest" not in models:
            print(f"  WARNING: Model '{MODEL}' not found in Ollama.")
            print(f"  Available: {', '.join(models) if models else 'none'}")
            print(f"  Run: ollama pull {MODEL}\n")
        else:
            print(f"  Ollama: connected. Model '{MODEL}' ready.\n")
    except:
        print("  WARNING: Cannot connect to Ollama at", OLLAMA_URL)
        print("  Start Ollama first. The UI will show 'offline' status.\n")

    # Load persistent memory
    global conversation_history
    loaded = load_chat_history()
    if loaded:
        conversation_history = loaded
        print(f"  Memory: loaded {len(loaded)} messages from previous sessions.")
    else:
        print(f"  Memory: fresh start (no previous history).")
    
    profile = load_user_profile()
    print(f"  Profile: {profile.get('name', 'Unknown')} | Last seen: {profile.get('last_seen', 'never')[:16]}")
    
    tasks = load_tasks()
    active = [t for t in tasks if t.get('status') != 'done']
    if active:
        print(f"  Tasks: {len(active)} active, {len(tasks) - len(active)} completed")
    
    # Start Telegram bot if configured
    try:
        if 'telegram_bot' in LOADED_SKILLS:
            tg_module = LOADED_SKILLS['telegram_bot']
            if hasattr(tg_module, 'start_telegram_bot'):
                tg_module.start_telegram_bot(
                    process_callback=process_message,
                    history_callback=conversation_history
                )
    except Exception as e:
        print(f"  Telegram: error starting â€” {e}")
    
    print()

    # Telegram-only mode: no web server, just Telegram polling
    if telegram_only:
        print("  Running in TELEGRAM-ONLY mode (no web UI).")
        print("  Send messages to the Telegram bot. Press Ctrl+C to stop.\n")
        try:
            while True:
                import time
                time.sleep(1)
        except KeyboardInterrupt:
            try:
                if 'telegram_bot' in LOADED_SKILLS:
                    LOADED_SKILLS['telegram_bot'].stop_telegram_bot()
            except:
                pass
            print("\n  Bye, Jane!")
            return

    # Open browser (only in normal mode)
    if not no_browser:
        def open_browser():
            import time
            time.sleep(1.0)
            webbrowser.open(f"http://localhost:{WEB_PORT}")
        threading.Thread(target=open_browser, daemon=True).start()

    # Start server
    server = HTTPServer(('0.0.0.0', WEB_PORT), AgentHandler)
    print(f"  Server running on http://localhost:{WEB_PORT}")
    if background_mode:
        print(f"  Background mode: web UI available but browser not opened.")
        print(f"  Send tasks from Telegram or open http://localhost:{WEB_PORT} manually.\n")
    else:
        print(f"  Press Ctrl+C to stop.\n")

    try:
        server.serve_forever()
    except KeyboardInterrupt:
        # Stop Telegram bot
        try:
            if 'telegram_bot' in LOADED_SKILLS:
                LOADED_SKILLS['telegram_bot'].stop_telegram_bot()
        except:
            pass
        print("\n  Shutting down...")
        server.server_close()
        print("  Bye, Jane!")


if __name__ == "__main__":
    main()

